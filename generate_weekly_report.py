"""
IPC Weekly Report Generator — Management Edition
Run: python3 generate_weekly_report.py
Output: IPC_Weekly_Report_<date>.pptx on Desktop

Week definition: Friday → Thursday
Sources:
  DAASH   → raw_dash.orders, branches, customers, revenueledgers
  GoSource→ raw_gosource.receipts
  AR      → gold.fact_ar_aging
"""

import os, datetime as dt
from urllib.parse import quote_plus
from dotenv import load_dotenv
import pandas as pd
from sqlalchemy import create_engine, text
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

load_dotenv()

# ── DB ────────────────────────────────────────────────────────────────────────
pw  = quote_plus(os.getenv("PG_PASSWORD", ""))
url = (
    f"postgresql+psycopg2://{os.getenv('PG_USER')}:{pw}"
    f"@{os.getenv('PG_HOST')}:{os.getenv('PG_PORT','25060')}"
    f"/PROD_ANALYTICS_DB?sslmode=require"
)
engine = create_engine(url, pool_pre_ping=True)

def q(sql):
    with engine.connect() as conn:
        return pd.read_sql(text(sql), conn)

# ── Dates (Week = Friday → Thursday) ─────────────────────────────────────────
today          = dt.date.today()
days_since_thu = (today.weekday() - 3) % 7      # Mon=0, Thu=3
tw_end         = today - dt.timedelta(days=days_since_thu)
tw_start       = tw_end   - dt.timedelta(days=6)
lw_end         = tw_start - dt.timedelta(days=1)
lw_start       = lw_end   - dt.timedelta(days=6)

TW_S, TW_E = tw_start.strftime("%Y-%m-%d"), tw_end.strftime("%Y-%m-%d")
LW_S, LW_E = lw_start.strftime("%Y-%m-%d"), lw_end.strftime("%Y-%m-%d")
TW_LABEL    = f"{tw_start.strftime('%d %b')} – {tw_end.strftime('%d %b %Y')}"
LW_LABEL    = f"{lw_start.strftime('%d %b')} – {lw_end.strftime('%d %b %Y')}"

print(f"TW: {TW_LABEL}   LW: {LW_LABEL}")
print("Fetching data...")

# ── GoSource: check if TW has delivered orders, else fall back to LW ──────────
_gs_tw_cnt = q(f"""
    SELECT COUNT(*) AS cnt FROM raw_gosource.receipts
    WHERE "createdAt"::date BETWEEN '{TW_S}' AND '{TW_E}'
      AND LOWER(status) = 'delivered'
""")
GS_TW_DELIVERED = int(_gs_tw_cnt.iloc[0]["cnt"])

if GS_TW_DELIVERED == 0:
    GS_TW_S, GS_TW_E = LW_S, LW_E
    GS_LW_S = (lw_start - dt.timedelta(days=7)).strftime("%Y-%m-%d")
    GS_LW_E = (lw_end   - dt.timedelta(days=7)).strftime("%Y-%m-%d")
    GS_TW_LABEL = LW_LABEL
    GS_LW_LABEL = (f"{(lw_start - dt.timedelta(days=7)).strftime('%d %b')} – "
                   f"{(lw_end   - dt.timedelta(days=7)).strftime('%d %b %Y')}")
    GS_NOTE = "  ·  0 delivered orders in current week — showing prior week"
else:
    GS_TW_S, GS_TW_E = TW_S, TW_E
    GS_LW_S, GS_LW_E = LW_S, LW_E
    GS_TW_LABEL, GS_LW_LABEL = TW_LABEL, LW_LABEL
    GS_NOTE = ""

# ══════════════════════════════════════════════════════════════════════════════
# DAASH QUERIES
# ══════════════════════════════════════════════════════════════════════════════

daash_kpi = q(f"""
    SELECT period,
        COUNT(*)  FILTER (WHERE status = 'delivered')                              AS orders,
        SUM("totalPrice") FILTER (WHERE status = 'delivered')                      AS revenue,
        SUM("totalPrice") FILTER (WHERE status = 'delivered')
            / NULLIF(COUNT(*) FILTER (WHERE status = 'delivered'), 0)              AS aov,
        COUNT(*) FILTER (WHERE status='delivered' AND channel='pos'
            AND LOWER(TRIM("paymentMethod")) NOT IN ('chowdeck','glovo'))          AS pos_orders,
        COUNT(*) FILTER (WHERE status='delivered' AND channel='website')           AS web_orders,
        COUNT(*) FILTER (WHERE status='delivered'
            AND LOWER(TRIM("paymentMethod"))='chowdeck')                           AS chow_orders,
        COUNT(*) FILTER (WHERE status='delivered'
            AND LOWER(TRIM("paymentMethod"))='glovo')                              AS glovo_orders,
        COUNT(*) FILTER (WHERE status = 'rejected')                                AS rejected,
        COUNT(*) FILTER (WHERE status = 'voided')                                  AS voided,
        COUNT(*)                                                                    AS total_all
    FROM (
        SELECT *, 'TW' AS period FROM raw_dash.orders
        WHERE "createdAt"::date BETWEEN '{TW_S}' AND '{TW_E}'
        UNION ALL
        SELECT *, 'LW' AS period FROM raw_dash.orders
        WHERE "createdAt"::date BETWEEN '{LW_S}' AND '{LW_E}'
    ) t
    GROUP BY period
""")

platform_fee = q(f"""
    SELECT period, SUM(amount) AS platform_fee
    FROM (
        SELECT 'TW' AS period, amount FROM raw_dash.revenueledgers
        WHERE type = 'credit'
          AND "createdAt"::date BETWEEN '{TW_S}' AND '{TW_E}'
        UNION ALL
        SELECT 'LW', amount FROM raw_dash.revenueledgers
        WHERE type = 'credit'
          AND "createdAt"::date BETWEEN '{LW_S}' AND '{LW_E}'
    ) t GROUP BY period
""")

# DOW: sort Fri→Thu  (postgres DOW: 0=Sun,1=Mon,2=Tue,3=Wed,4=Thu,5=Fri,6=Sat)
daash_dow = q(f"""
    SELECT
        EXTRACT(DOW FROM "createdAt")::int           AS dow_num,
        TO_CHAR(MIN("createdAt"), 'Dy')              AS day_name,
        SUM(CASE WHEN period='TW' THEN 1 ELSE 0 END) AS tw_orders,
        SUM(CASE WHEN period='LW' THEN 1 ELSE 0 END) AS lw_orders
    FROM (
        SELECT "createdAt", 'TW' AS period FROM raw_dash.orders
        WHERE "createdAt"::date BETWEEN '{TW_S}' AND '{TW_E}' AND status='delivered'
        UNION ALL
        SELECT "createdAt", 'LW' AS period FROM raw_dash.orders
        WHERE "createdAt"::date BETWEEN '{LW_S}' AND '{LW_E}' AND status='delivered'
    ) t
    GROUP BY EXTRACT(DOW FROM "createdAt")
    ORDER BY (EXTRACT(DOW FROM "createdAt")::int - 5 + 7) % 7
""")

# Brand order counts (no revenue)
daash_brands_tw = q(f"""
    SELECT
        COALESCE(NULLIF(TRIM(c."businessName"), ''), 'Unknown') AS brand,
        COUNT(*) AS orders_tw
    FROM raw_dash.orders o
    JOIN raw_dash.branches b ON o.branch = b._id
    JOIN raw_dash.customers c ON b."customerId" = c._id
    WHERE o."createdAt"::date BETWEEN '{TW_S}' AND '{TW_E}'
      AND o.status = 'delivered'
    GROUP BY c."businessName"
    ORDER BY orders_tw DESC
    LIMIT 8
""")
daash_brands_lw = q(f"""
    SELECT
        COALESCE(NULLIF(TRIM(c."businessName"), ''), 'Unknown') AS brand,
        COUNT(*) AS orders_lw
    FROM raw_dash.orders o
    JOIN raw_dash.branches b ON o.branch = b._id
    JOIN raw_dash.customers c ON b."customerId" = c._id
    WHERE o."createdAt"::date BETWEEN '{LW_S}' AND '{LW_E}'
      AND o.status = 'delivered'
    GROUP BY c."businessName"
""")
daash_brands = daash_brands_tw.merge(daash_brands_lw, on="brand", how="left")
daash_brands["orders_lw"] = daash_brands["orders_lw"].fillna(0)

# Brand retention
daash_retention = q(f"""
    WITH tw_b AS (
        SELECT DISTINCT b."customerId" AS bid FROM raw_dash.orders o
        JOIN raw_dash.branches b ON o.branch = b._id
        WHERE o."createdAt"::date BETWEEN '{TW_S}' AND '{TW_E}' AND o.status='delivered'
    ),
    lw_b AS (
        SELECT DISTINCT b."customerId" AS bid FROM raw_dash.orders o
        JOIN raw_dash.branches b ON o.branch = b._id
        WHERE o."createdAt"::date BETWEEN '{LW_S}' AND '{LW_E}' AND o.status='delivered'
    ),
    ever_b AS (
        SELECT DISTINCT b."customerId" AS bid FROM raw_dash.orders o
        JOIN raw_dash.branches b ON o.branch = b._id
        WHERE o."createdAt"::date < '{TW_S}' AND o.status='delivered'
    )
    SELECT
        (SELECT COUNT(*) FROM tw_b)                                          AS active_tw,
        (SELECT COUNT(*) FROM lw_b)                                          AS active_lw,
        (SELECT COUNT(*) FROM tw_b WHERE bid IN (SELECT bid FROM lw_b))      AS retained,
        (SELECT COUNT(*) FROM tw_b WHERE bid NOT IN (SELECT bid FROM ever_b)) AS new_brands,
        (SELECT COUNT(*) FROM lw_b WHERE bid NOT IN (SELECT bid FROM tw_b))  AS churned
""")

daash_issues_by_brand = q(f"""
    SELECT
        COALESCE(NULLIF(TRIM(c."businessName"), ''), 'Unknown') AS brand,
        COUNT(*) AS total_orders,
        COUNT(*) FILTER (WHERE o.status = 'rejected') AS rejected,
        COUNT(*) FILTER (WHERE o.status = 'voided')   AS voided,
        ROUND(COUNT(*) FILTER (WHERE o.status IN ('rejected','voided'))
            * 100.0 / NULLIF(COUNT(*), 0), 1) AS issue_rate
    FROM raw_dash.orders o
    JOIN raw_dash.branches b ON o.branch = b._id
    JOIN raw_dash.customers c ON b."customerId" = c._id
    WHERE o."createdAt"::date BETWEEN '{TW_S}' AND '{TW_E}'
      AND o.status IN ('delivered', 'rejected', 'voided')
    GROUP BY c."businessName"
    HAVING COUNT(*) >= 5
    ORDER BY issue_rate DESC
    LIMIT 6
""")

# ══════════════════════════════════════════════════════════════════════════════
# GOSOURCE QUERIES
# ══════════════════════════════════════════════════════════════════════════════

gs_kpi = q(f"""
    SELECT period,
        SUM("totalPrice")                    AS revenue,
        SUM("totalPrice")/NULLIF(COUNT(*),0) AS aov
    FROM (
        SELECT *, 'TW' AS period FROM raw_gosource.receipts
        WHERE "createdAt"::date BETWEEN '{GS_TW_S}' AND '{GS_TW_E}'
          AND LOWER(status)='delivered'
        UNION ALL
        SELECT *, 'LW' AS period FROM raw_gosource.receipts
        WHERE "createdAt"::date BETWEEN '{GS_LW_S}' AND '{GS_LW_E}'
          AND LOWER(status)='delivered'
    ) t GROUP BY period
""")

# Total orders placed (any status) + active customers (any status)
gs_activity = q(f"""
    SELECT period,
        COUNT(*)                     AS orders,
        COUNT(DISTINCT business)     AS customers
    FROM (
        SELECT business, 'TW' AS period FROM raw_gosource.receipts
        WHERE "createdAt"::date BETWEEN '{GS_TW_S}' AND '{GS_TW_E}'
        UNION ALL
        SELECT business, 'LW' AS period FROM raw_gosource.receipts
        WHERE "createdAt"::date BETWEEN '{GS_LW_S}' AND '{GS_LW_E}'
    ) t GROUP BY period
""")

gs_payment = q(f"""
    SELECT period,
        COUNT(*) FILTER (WHERE LOWER("paymentMethod")='credit')               AS credit_orders,
        COUNT(*) FILTER (WHERE LOWER(COALESCE("paymentMethod",''))!='credit') AS other_orders,
        COUNT(*) AS total_orders
    FROM (
        SELECT "paymentMethod", 'TW' AS period FROM raw_gosource.receipts
        WHERE "createdAt"::date BETWEEN '{GS_TW_S}' AND '{GS_TW_E}'
        UNION ALL
        SELECT "paymentMethod", 'LW' AS period FROM raw_gosource.receipts
        WHERE "createdAt"::date BETWEEN '{GS_LW_S}' AND '{GS_LW_E}'
    ) t GROUP BY period
""")

gs_dow = q(f"""
    SELECT
        EXTRACT(DOW FROM "createdAt")::int AS dow_num,
        TO_CHAR(MIN("createdAt"), 'Dy')    AS day_name,
        COUNT(*) AS orders
    FROM raw_gosource.receipts
    WHERE "createdAt"::date BETWEEN '{GS_TW_S}' AND '{GS_TW_E}'
      AND LOWER(status)='delivered'
    GROUP BY EXTRACT(DOW FROM "createdAt")
    ORDER BY (EXTRACT(DOW FROM "createdAt")::int - 5 + 7) % 7
""")

gs_pipeline = q("""
    SELECT INITCAP(LOWER(status)) AS status, COUNT(*) AS orders
    FROM raw_gosource.receipts
    WHERE "createdAt"::date >= CURRENT_DATE - INTERVAL '30 days'
    GROUP BY LOWER(status)
    ORDER BY orders DESC
""")

gs_ar = q("""
    SELECT
        CASE
            WHEN CURRENT_DATE - "createdAt"::date <= 30 THEN '0-30 days'
            WHEN CURRENT_DATE - "createdAt"::date <= 60 THEN '31-60 days'
            WHEN CURRENT_DATE - "createdAt"::date <= 90 THEN '61-90 days'
            ELSE '90+ days'
        END AS ar_aging_bucket,
        COUNT(*)            AS invoices,
        SUM("totalPrice")   AS amount
    FROM raw_gosource.receipts
    WHERE LOWER("paymentMethod") = 'credit'
      AND LOWER(status) = 'delivered'
      AND LOWER(COALESCE("paymentStatus", '')) != 'paid'
    GROUP BY ar_aging_bucket
    ORDER BY ar_aging_bucket
""")

print("Data fetched. Building slides...")

# ── Scalar helpers ─────────────────────────────────────────────────────────────
def _get(df, period, col):
    row = df[df["period"] == period]
    if row.empty or pd.isna(row.iloc[0][col]): return 0.0
    return float(row.iloc[0][col])

def _get_pf(period):
    row = platform_fee[platform_fee["period"] == period]
    return float(row.iloc[0]["platform_fee"] or 0) if not row.empty else 0.0

def naira(v):
    if v >= 1_000_000_000: return f"₦{v/1_000_000_000:.1f}B"
    if v >= 1_000_000:     return f"₦{v/1_000_000:.1f}M"
    if v >= 1_000:         return f"₦{v/1_000:.0f}K"
    return f"₦{v:.0f}"

def chg(c, p):
    if p and p > 0: return (c - p) / p * 100
    return None

def arrow(curr, prev, good="up"):
    pct = chg(curr, prev)
    if pct is None: return "— vs prior wk", SLATE
    sym   = "▲" if pct > 0 else ("▼" if pct < 0 else "→")
    color = _trend_color(pct, good)
    return f"{sym} {abs(pct):.1f}% vs prior wk", color

def _trend_color(pct, good="up"):
    if pct is None: return SLATE
    if good == "up":
        return GREEN if pct > 0 else (AMBER if pct > -5 else RED)
    return GREEN if pct < 0 else (AMBER if pct < 5 else RED)

# DAASH scalars
d_tw_ord   = _get(daash_kpi, "TW", "orders");    d_lw_ord  = _get(daash_kpi, "LW", "orders")
d_tw_rev   = _get(daash_kpi, "TW", "revenue");   d_lw_rev  = _get(daash_kpi, "LW", "revenue")
d_tw_aov   = _get(daash_kpi, "TW", "aov");       d_lw_aov  = _get(daash_kpi, "LW", "aov")
d_tw_svc   = _get_pf("TW");                      d_lw_svc  = _get_pf("LW")
d_tw_pos   = _get(daash_kpi, "TW", "pos_orders"); d_lw_pos  = _get(daash_kpi, "LW", "pos_orders")
d_tw_web   = _get(daash_kpi, "TW", "web_orders"); d_lw_web  = _get(daash_kpi, "LW", "web_orders")
d_tw_chow  = _get(daash_kpi, "TW", "chow_orders");d_lw_chow = _get(daash_kpi, "LW", "chow_orders")
d_tw_glov  = _get(daash_kpi, "TW", "glovo_orders");d_lw_glov= _get(daash_kpi, "LW", "glovo_orders")
d_tw_rej   = _get(daash_kpi, "TW", "rejected");  d_lw_rej  = _get(daash_kpi, "LW", "rejected")
d_tw_void  = _get(daash_kpi, "TW", "voided");    d_lw_void = _get(daash_kpi, "LW", "voided")
d_tw_total = _get(daash_kpi, "TW", "total_all"); d_lw_total= _get(daash_kpi, "LW", "total_all")
d_tw_iss   = d_tw_rej + d_tw_void
d_lw_iss   = d_lw_rej + d_lw_void
d_tw_irate = d_tw_iss / d_tw_total * 100 if d_tw_total > 0 else 0
d_lw_irate = d_lw_iss / d_lw_total * 100 if d_lw_total > 0 else 0
d_tw_comp  = d_tw_ord / d_tw_total * 100 if d_tw_total > 0 else 0

# Brand retention scalars
ret = daash_retention.iloc[0] if not daash_retention.empty else None
r_active_tw  = int(ret["active_tw"])  if ret is not None else 0
r_active_lw  = int(ret["active_lw"])  if ret is not None else 0
r_retained   = int(ret["retained"])   if ret is not None else 0
r_new        = int(ret["new_brands"]) if ret is not None else 0
r_churned    = int(ret["churned"])    if ret is not None else 0

# GoSource scalars
g_tw_rev  = _get(gs_kpi, "TW", "revenue"); g_lw_rev  = _get(gs_kpi, "LW", "revenue")
g_tw_aov  = _get(gs_kpi, "TW", "aov");    g_lw_aov  = _get(gs_kpi, "LW", "aov")
g_tw_ord  = _get(gs_activity, "TW", "orders");  g_lw_ord  = _get(gs_activity, "LW", "orders")
g_tw_cust = _get(gs_activity, "TW", "customers"); g_lw_cust = _get(gs_activity, "LW", "customers")

def _gs_pay(period, col):
    row = gs_payment[gs_payment["period"] == period]
    return float(row.iloc[0][col]) if not row.empty else 0.0

g_tw_credit = _gs_pay("TW", "credit_orders"); g_tw_cash = _gs_pay("TW", "cash_orders")
g_tw_pay_tot = g_tw_credit + g_tw_cash
g_tw_credit_pct = g_tw_credit / g_tw_pay_tot * 100 if g_tw_pay_tot > 0 else 0
g_tw_cash_pct   = g_tw_cash   / g_tw_pay_tot * 100 if g_tw_pay_tot > 0 else 0

ar_total = float(gs_ar["amount"].sum())      if not gs_ar.empty else 0
ar_90    = float(gs_ar[gs_ar["ar_aging_bucket"]=="90+ days"]["amount"].sum()) if not gs_ar.empty else 0
ar_90pct = ar_90 / ar_total * 100 if ar_total > 0 else 0

print(f"  DAASH TW: {int(d_tw_ord):,} orders  ·  issue rate {d_tw_irate:.2f}%  ·  svc fee {naira(d_tw_svc)}")
print(f"  GoSrc TW: {int(g_tw_ord):,} orders  ·  {naira(g_tw_rev)}  ·  credit {g_tw_credit_pct:.0f}%")

# ── Colour palette ─────────────────────────────────────────────────────────────
DAASH_C  = RGBColor(0x7F, 0x1D, 0x1D)   # deep crimson
GS_C     = RGBColor(0x1A, 0x4D, 0x2E)   # forest green
DARK     = RGBColor(0x0F, 0x17, 0x2A)   # near-black navy
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE= RGBColor(0xF9, 0xFA, 0xFB)
LIGHT_BG = RGBColor(0xF3, 0xF4, 0xF6)
MID_GRAY = RGBColor(0x6B, 0x72, 0x80)
SLATE    = RGBColor(0x9C, 0xA3, 0xAF)
BORDER   = RGBColor(0xE5, 0xE7, 0xEB)
GREEN    = RGBColor(0x16, 0xA3, 0x4A)
AMBER    = RGBColor(0xD9, 0x77, 0x06)
RED      = RGBColor(0xDC, 0x26, 0x26)
BLUE     = RGBColor(0x25, 0x63, 0xEB)
PURPLE   = RGBColor(0x7C, 0x3A, 0xED)
ORANGE   = RGBColor(0xEA, 0x58, 0x0C)
TEAL     = RGBColor(0x0D, 0x94, 0x88)

# ── Presentation setup ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def new_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

# ── Primitive drawing ─────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill=WHITE, border=None, border_pt=0.5):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(border_pt)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, size=11, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    return txb

# ── Higher-level components ───────────────────────────────────────────────────
def page_header(slide, title, sub, accent):
    """Full-width clean header: colored left accent bar + title/sub text."""
    rect(slide, 0, 0, 13.33, 1.05, fill=accent)
    rect(slide, 0, 0, 0.5,   1.05, fill=_darken(accent))   # left accent stripe
    txt(slide, title, 0.65, 0.08, 12.5, 0.56, size=20, bold=True,  color=WHITE)
    txt(slide, sub,   0.65, 0.67, 12.5, 0.32, size=9.5, color=RGBColor(0xD1,0xD5,0xDB))

def _darken(c):
    """Make colour ~25% darker for accent stripe."""
    def d(v): return max(0, v - 60)
    return RGBColor(d(c[0]), d(c[1]), d(c[2]))

def kpi_card(slide, x, y, w, h, label, value,
             arrow_txt=None, arrow_col=None, val_size=30, accent=None):
    """Clean KPI card: white bg, thin border, top-left accent dot."""
    rect(slide, x, y, w, h, fill=WHITE, border=BORDER, border_pt=0.8)
    if accent:
        rect(slide, x, y, w, 0.055, fill=accent)   # thin top accent line
    pad = 0.2
    txt(slide, label.upper(), x+pad, y+0.17, w-pad-0.08, 0.28,
        size=8, bold=True, color=SLATE)
    txt(slide, str(value), x+pad, y+0.44, w-pad-0.08, h*0.45,
        size=val_size, bold=True, color=DARK, wrap=False)
    if arrow_txt:
        txt(slide, arrow_txt, x+pad, y+h-0.38, w-pad-0.08, 0.3,
            size=9, bold=True, color=arrow_col or SLATE, wrap=False)

def insight_box(slide, text, x, y, w, h, bg_col=None, border_col=None,
                text_col=DARK, size=10):
    bg_col     = bg_col     or RGBColor(0xEF,0xF6,0xFF)
    border_col = border_col or BLUE
    rect(slide, x, y, w, h, fill=bg_col, border=border_col, border_pt=1.2)
    rect(slide, x, y, 0.045, h, fill=border_col)       # left accent stripe
    txt(slide, text, x+0.22, y+0.14, w-0.35, h-0.24,
        size=size, color=text_col, wrap=True)

def section_label(slide, text, x, y, w, accent):
    """Small pill-style section separator."""
    rect(slide, x, y, w, 0.28, fill=accent)
    txt(slide, text, x+0.15, y+0.03, w-0.2, 0.22,
        size=8.5, bold=True, color=WHITE)

def divider(slide, x, y, h):
    rect(slide, x, y, 0.03, h, fill=BORDER)

def table_block(slide, headers, rows, x, y, col_widths, row_h=0.38, header_color=None):
    header_color = header_color or DARK
    cx = x
    for i, h in enumerate(headers):
        rect(slide, cx, y, col_widths[i], row_h, fill=header_color)
        tf = slide.shapes.add_textbox(
            Inches(cx+0.08), Inches(y+0.06),
            Inches(col_widths[i]-0.12), Inches(row_h-0.08)).text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(h)
        r.font.size = Pt(8.5); r.font.bold = True; r.font.color.rgb = WHITE
        cx += col_widths[i]
    for ri, row in enumerate(rows):
        cx = x
        row_bg = WHITE if ri % 2 == 0 else OFF_WHITE
        for ci, (cell, cw) in enumerate(zip(row, col_widths)):
            is_colored = isinstance(cell, tuple)
            cell_txt   = cell[0] if is_colored else str(cell)
            cell_col   = cell[1] if is_colored else DARK
            rect(slide, cx, y+row_h*(ri+1), cw, row_h,
                 fill=row_bg, border=BORDER, border_pt=0.3)
            tf = slide.shapes.add_textbox(
                Inches(cx+0.08), Inches(y+row_h*(ri+1)+0.06),
                Inches(cw-0.1),  Inches(row_h-0.09)).text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(cell_txt)
            r.font.size = Pt(9.5); r.font.color.rgb = cell_col
            if ci == 0: r.font.bold = True
            cx += cw

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 0 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, DARK)

# Left sidebar accent
rect(sl, 0,    0, 0.55, 7.5, fill=DAASH_C)
rect(sl, 0.55, 0, 0.18, 7.5, fill=GS_C)

# Title block
txt(sl, "IPC GROUP", 1.05, 1.6, 11.5, 0.55, size=12, bold=True,
    color=RGBColor(0x6B,0x72,0x80))
txt(sl, "Weekly Performance", 1.05, 2.12, 11.5, 0.9, size=42, bold=True, color=WHITE)
txt(sl, "Report", 1.05, 2.96, 11.5, 0.9, size=42, bold=True,
    color=RGBColor(0xD1,0xD5,0xDB))

# Week line
rect(sl, 1.05, 4.05, 8.5, 0.05, fill=RGBColor(0x37,0x41,0x51))
txt(sl, f"Week  {TW_LABEL}   ·   vs Prior  {LW_LABEL}",
    1.05, 4.2, 11.5, 0.45, size=13, color=RGBColor(0x9C,0xA3,0xAF))
txt(sl, f"Prepared  {today.strftime('%A, %d %B %Y')}",
    1.05, 4.72, 11.5, 0.38, size=10, color=RGBColor(0x6B,0x72,0x80))

# Service badges
for bx, label, col in [(1.05, "DAASH  ·  Food Delivery", DAASH_C),
                        (5.35, "GOSOURCE  ·  B2B Procurement", GS_C)]:
    rect(sl, bx, 5.6, 3.9, 0.6, fill=RGBColor(0x1F,0x29,0x3B), border=col, border_pt=1.2)
    txt(sl, label, bx+0.18, 5.68, 3.6, 0.42, size=11, bold=True, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — EXECUTIVE SUMMARY
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, LIGHT_BG)
page_header(sl, "EXECUTIVE SUMMARY",
            f"DAASH: {TW_LABEL}   ·   GoSource: {GS_TW_LABEL}{GS_NOTE}", DARK)

# DAASH column
rect(sl, 0.25, 1.12, 6.1, 0.3, fill=DAASH_C)
txt(sl, "DAASH  —  FOOD DELIVERY", 0.4, 1.15, 5.9, 0.24, size=9, bold=True, color=WHITE)

KW, KH, KG = 2.92, 1.52, 0.13
for i, (lbl, val, p1, p2) in enumerate([
    ("Total Orders",    f"{int(d_tw_ord):,}",   d_tw_ord, d_lw_ord),
    ("Service Fee",     naira(d_tw_svc),         d_tw_svc, d_lw_svc),
    ("Avg Order Value", naira(d_tw_aov),         d_tw_aov, d_lw_aov),
    ("Issue Rate",      f"{d_tw_irate:.2f}%",    d_lw_irate, d_tw_irate),  # lower=better
]):
    col = i % 2; row = i // 2
    tx = 0.25 + col*(KW+KG)
    ty = 1.48 + row*(KH+KG)
    g = "up" if lbl != "Issue Rate" else "down"
    at, ac = arrow(p1, p2, good=g)
    kpi_card(sl, tx, ty, KW, KH, lbl, val, at, ac, val_size=22, accent=DAASH_C)

# GoSource column
rect(sl, 6.98, 1.12, 6.1, 0.3, fill=GS_C)
txt(sl, f"GOSOURCE  —  B2B PROCUREMENT  ({GS_TW_LABEL})",
    7.12, 1.15, 5.9, 0.24, size=9, bold=True, color=WHITE)

for i, (lbl, val, p1, p2) in enumerate([
    ("Total Orders",     f"{int(g_tw_ord):,}",   g_tw_ord,  g_lw_ord),
    ("Avg Order Value",  naira(g_tw_aov),         g_tw_aov,  g_lw_aov),
    ("Active Customers", f"{int(g_tw_cust):,}",  g_tw_cust, g_lw_cust),
    ("Credit Orders",    f"{g_tw_credit_pct:.0f}%", g_tw_credit_pct, None),
]):
    col = i % 2; row = i // 2
    tx = 6.98 + col*(KW+KG)
    ty = 1.48 + row*(KH+KG)
    at, ac = arrow(p1, p2) if p2 else ("— B2B mix", SLATE)
    kpi_card(sl, tx, ty, KW, KH, lbl, val, at, ac, val_size=22, accent=GS_C)

divider(sl, 6.62, 1.12, 3.52)

# Alert row
irate_bg = RGBColor(0xFE,0xF2,0xF2) if d_tw_irate > 1 else RGBColor(0xF0,0xFD,0xF4)
irate_bd = RED if d_tw_irate > 1 else GREEN
daash_alert = (f"Issue rate {d_tw_irate:.2f}%  ·  "
               f"{int(d_tw_rej):,} rejected + {int(d_tw_void):,} voided  ·  "
               f"Completion rate {d_tw_comp:.1f}%")
insight_box(sl, daash_alert, 0.25, 5.18, 6.1, 0.65,
            bg_col=irate_bg, border_col=irate_bd,
            text_col=RGBColor(0x99,0x1B,0x1B) if d_tw_irate > 1 else RGBColor(0x14,0x53,0x2D))

ar_bg = RGBColor(0xFE,0xF2,0xF2) if ar_90pct > 15 else RGBColor(0xF0,0xFD,0xF4)
ar_bd = RED if ar_90pct > 15 else GREEN
ar_msg = (f"{naira(ar_90)} ({ar_90pct:.1f}% of AR) overdue 90+ days — action needed"
          if ar_90pct > 0 else f"AR healthy — {naira(ar_total)} total outstanding, no overdue concentration")
insight_box(sl, ar_msg, 6.98, 5.18, 6.1, 0.65,
            bg_col=ar_bg, border_col=ar_bd,
            text_col=RGBColor(0x99,0x1B,0x1B) if ar_90pct > 15 else RGBColor(0x14,0x53,0x2D))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — DAASH: ORDERS & PERFORMANCE
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, LIGHT_BG)
page_header(sl, "DAASH  —  ORDERS & PERFORMANCE",
            f"Week: {TW_LABEL}   ·   vs Prior: {LW_LABEL}", DAASH_C)

# 4 main KPI tiles (row 1)
KW2, KH2, KG2 = 3.03, 2.1, 0.12
for i, (lbl, val, p1, p2, g) in enumerate([
    ("Total Orders",    f"{int(d_tw_ord):,}", d_tw_ord, d_lw_ord, "up"),
    ("Avg Order Value", naira(d_tw_aov),      d_tw_aov, d_lw_aov, "up"),
    ("Platform Fee",    naira(d_tw_svc),      d_tw_svc, d_lw_svc, "up"),
    ("Completion Rate", f"{d_tw_comp:.1f}%",  d_tw_comp,
     (d_lw_ord / d_lw_total * 100 if d_lw_total > 0 else 0), "up"),
]):
    at, ac = arrow(p1, p2, good=g)
    kpi_card(sl, 0.28 + i*(KW2+KG2), 1.22, KW2, KH2, lbl, val, at, ac,
             val_size=34, accent=DAASH_C)

# Insight
rev_chg = chg(d_tw_ord, d_lw_ord)
svc_chg = chg(d_tw_svc, d_lw_svc)
ins2 = (f"{'▲' if (rev_chg or 0)>0 else '▼'} "
        f"{abs(rev_chg or 0):.1f}% WoW on delivered orders — "
        f"{int(d_tw_ord):,} orders at {naira(d_tw_aov)} avg. "
        f"Platform service fee: {naira(d_tw_svc)} "
        f"({'up' if (svc_chg or 0)>0 else 'down'} {abs(svc_chg or 0):.1f}% WoW). "
        f"Completion rate {d_tw_comp:.1f}% — "
        f"{'strong' if d_tw_comp > 95 else 'watch the issue rate closely'}.")
insight_box(sl, ins2, 0.28, 3.48, 12.77, 0.7)

# Retention quick stats (bottom)
rect(sl, 0.28, 4.32, 12.77, 0.28, fill=DAASH_C)
txt(sl, "BRAND ACTIVITY THIS WEEK", 0.45, 4.35, 12.5, 0.22,
    size=8.5, bold=True, color=WHITE)

stat_items = [
    ("Active Brands", f"{r_active_tw}", DARK),
    ("Retained",      f"{r_retained}", GREEN),
    ("New This Wk",   f"{r_new}", BLUE),
    ("Not Ordering",  f"{r_churned}", RED),
    ("Prior Week",    f"{r_active_lw}", MID_GRAY),
]
SW = 12.77 / len(stat_items)
for i, (lbl, val, col) in enumerate(stat_items):
    bx = 0.28 + i*SW
    bg_stat = RGBColor(0xF9,0xFA,0xFB) if i % 2 == 0 else WHITE
    rect(sl, bx, 4.65, SW, 1.65, fill=bg_stat, border=BORDER, border_pt=0.5)
    txt(sl, lbl, bx+0.12, 4.78, SW-0.18, 0.28, size=8, color=SLATE, bold=True)
    txt(sl, val,  bx+0.12, 5.05, SW-0.18, 0.9, size=30, bold=True, color=col, wrap=False)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — DAASH: CHANNEL BREAKDOWN
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, LIGHT_BG)
page_header(sl, "DAASH  —  CHANNEL BREAKDOWN",
            f"Week: {TW_LABEL}   ·   Orders by channel — no aggregation grouping", DAASH_C)

CH_DATA = [
    ("POS",      d_tw_pos,  d_lw_pos,  DAASH_C),
    ("Website",  d_tw_web,  d_lw_web,  BLUE),
    ("Chowdeck", d_tw_chow, d_lw_chow, ORANGE),
    ("Glovo",    d_tw_glov, d_lw_glov, TEAL),
]

CW3, CH3, CG3 = 3.0, 2.6, 0.12
for i, (lbl, o_tw, o_lw, acc) in enumerate(CH_DATA):
    tx = 0.28 + i*(CW3+CG3)
    pct = o_tw / d_tw_ord * 100 if d_tw_ord > 0 else 0
    at, ac = arrow(o_tw, o_lw)

    # Card
    rect(sl, tx, 1.22, CW3, CH3, fill=WHITE, border=BORDER, border_pt=0.8)
    rect(sl, tx, 1.22, CW3, 0.055, fill=acc)   # top strip
    txt(sl, lbl.upper(), tx+0.2, 1.38, CW3-0.3, 0.28, size=9, bold=True, color=acc)
    txt(sl, f"{int(o_tw):,}", tx+0.2, 1.66, CW3-0.3, 0.75, size=36, bold=True, color=DARK, wrap=False)
    txt(sl, "orders", tx+0.2, 2.38, CW3-0.3, 0.3, size=10, color=MID_GRAY)
    # % share bar
    bar_w = (CW3 - 0.4) * pct / 100
    rect(sl, tx+0.2, 2.72, CW3-0.4, 0.16, fill=LIGHT_BG, border=BORDER, border_pt=0.3)
    if bar_w > 0:
        rect(sl, tx+0.2, 2.72, bar_w, 0.16, fill=acc)
    txt(sl, f"{pct:.1f}% of total orders", tx+0.2, 2.9, CW3-0.3, 0.28, size=9, color=MID_GRAY)
    txt(sl, at, tx+0.2, 3.22, CW3-0.3, 0.28, size=9, bold=True, color=ac, wrap=False)

# Insight
pos_pct   = d_tw_pos   / d_tw_ord * 100 if d_tw_ord > 0 else 0
web_pct   = d_tw_web   / d_tw_ord * 100 if d_tw_ord > 0 else 0
chow_pct  = d_tw_chow  / d_tw_ord * 100 if d_tw_ord > 0 else 0
glov_pct  = d_tw_glov  / d_tw_ord * 100 if d_tw_ord > 0 else 0
agg_pct   = chow_pct + glov_pct

ins3 = (f"POS {pos_pct:.0f}%  ·  Website {web_pct:.0f}%  ·  "
        f"Chowdeck {chow_pct:.0f}%  ·  Glovo {glov_pct:.0f}% of total volume. "
        f"Aggregator channels (Chowdeck + Glovo) account for {agg_pct:.0f}% of orders — "
        f"these carry third-party commissions that reduce net margin. "
        f"Target: grow POS and Website share week-on-week.")
insight_box(sl, ins3, 0.28, 4.02, 12.77, 0.72)

# Channel comparison table
rows3 = []
for lbl, o_tw, o_lw, acc in CH_DATA:
    pct_t = o_tw / d_tw_ord * 100 if d_tw_ord > 0 else 0
    pct_l = o_lw / d_lw_ord * 100 if d_lw_ord > 0 else 0
    chg_o = chg(o_tw, o_lw)
    chg_sym = (f"▲ {abs(chg_o):.1f}%" if (chg_o or 0)>0 else f"▼ {abs(chg_o or 0):.1f}%")
    chg_col = GREEN if (chg_o or 0) > 0 else RED
    rows3.append((
        lbl,
        f"{int(o_tw):,}",
        f"{pct_t:.1f}%",
        f"{int(o_lw):,}",
        f"{pct_l:.1f}%",
        (chg_sym, chg_col),
    ))
table_block(sl, ["Channel", "TW Orders", "TW Share %", "LW Orders", "LW Share %", "Change"],
            rows3, x=0.28, y=4.87, col_widths=[2.2, 2.0, 2.0, 2.0, 2.1, 2.47],
            row_h=0.36, header_color=DAASH_C)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — DAASH: SALES BY DAY OF WEEK
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, LIGHT_BG)
page_header(sl, "DAASH  —  SALES BY DAY OF WEEK",
            f"Week: {TW_LABEL}   ·   vs Prior: {LW_LABEL}   ·   Week starts Friday", DAASH_C)

if not daash_dow.empty:
    max_tw = int(daash_dow["tw_orders"].max()) or 1

    BAR_TOP    = 1.3
    BAR_AREA_H = 3.8
    LEGEND_Y   = 5.2

    col_count = len(daash_dow)
    col_w     = 12.77 / col_count
    bar_max_h = BAR_AREA_H - 0.7

    for i, row in daash_dow.iterrows():
        dow  = int(row["dow_num"])
        name = str(row["day_name"])
        tw_o = int(row["tw_orders"])
        lw_o = int(row["lw_orders"])
        col_x = 0.28 + i * col_w

        # Background alternating
        alt = RGBColor(0xF8,0xF9,0xFA) if i % 2 == 0 else WHITE
        rect(sl, col_x, BAR_TOP, col_w, BAR_AREA_H + 0.9, fill=alt, border=BORDER, border_pt=0.3)

        # Day label
        txt(sl, name, col_x, BAR_TOP + 0.05, col_w, 0.3,
            size=9, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

        # LW bar (background, lighter)
        lw_h = (lw_o / max_tw) * bar_max_h
        lw_y  = BAR_TOP + BAR_AREA_H - lw_h - 0.35
        rect(sl, col_x + col_w*0.35, lw_y, col_w * 0.22, lw_h,
             fill=RGBColor(0xD1,0xD5,0xDB))

        # TW bar
        tw_h = (tw_o / max_tw) * bar_max_h
        tw_y  = BAR_TOP + BAR_AREA_H - tw_h - 0.35
        rect(sl, col_x + col_w*0.12, tw_y, col_w * 0.22, tw_h, fill=DAASH_C)

        # TW count label
        txt(sl, f"{tw_o:,}", col_x, BAR_TOP + BAR_AREA_H - 0.28, col_w, 0.26,
            size=8.5, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # Legend
    rect(sl, 0.28, LEGEND_Y, 0.18, 0.18, fill=DAASH_C)
    txt(sl, "This Week", 0.52, LEGEND_Y, 1.5, 0.2, size=8.5, color=DARK)
    rect(sl, 2.1, LEGEND_Y, 0.18, 0.18, fill=RGBColor(0xD1,0xD5,0xDB))
    txt(sl, "Prior Week", 2.34, LEGEND_Y, 1.5, 0.2, size=8.5, color=DARK)

    # Peak day insight
    peak_row = daash_dow.loc[daash_dow["tw_orders"].idxmax()]
    slow_row = daash_dow.loc[daash_dow["tw_orders"].idxmin()]
    ins4 = (f"Peak day: {peak_row['day_name']} with {int(peak_row['tw_orders']):,} orders. "
            f"Slowest: {slow_row['day_name']} with {int(slow_row['tw_orders']):,} orders. "
            f"Use this pattern for staffing, rider allocation, and kitchen prep scheduling. "
            f"If weekday peaks are growing, B2B lunch orders may be driving volume.")
    insight_box(sl, ins4, 0.28, 5.52, 12.77, 0.72)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DAASH: BRAND ACTIVITY
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, LIGHT_BG)
page_header(sl, "DAASH  —  BRAND ACTIVITY & ORDER SHARE",
            f"Week: {TW_LABEL}   ·   Order counts & share — revenue not shown", DAASH_C)

# Retention stat row
ret_items = [
    ("Active Brands TW", str(r_active_tw), DAASH_C),
    ("Active Brands LW", str(r_active_lw), MID_GRAY),
    ("Retained",         str(r_retained),  GREEN),
    ("New This Week",    str(r_new),        BLUE),
    ("Not Ordering",     str(r_churned),    RED),
]
RW = 12.77 / len(ret_items)
for i, (lbl, val, col) in enumerate(ret_items):
    rx = 0.28 + i*RW
    rect(sl, rx, 1.15, RW, 1.28, fill=WHITE, border=BORDER, border_pt=0.8)
    rect(sl, rx, 1.15, RW, 0.055, fill=col)
    txt(sl, lbl.upper(), rx+0.15, 1.3, RW-0.2, 0.26, size=7.5, bold=True, color=SLATE)
    txt(sl, val, rx+0.15, 1.55, RW-0.2, 0.68, size=26, bold=True, color=col, wrap=False)

# Brand table (order count + % share, NO revenue)
if not daash_brands.empty:
    total_tw_ord = daash_brands["orders_tw"].sum() or 1
    rows5 = []
    for rk, br in enumerate(daash_brands.itertuples(), 1):
        share_tw = br.orders_tw / total_tw_ord * 100
        share_lw = br.orders_lw / daash_brands["orders_lw"].sum() * 100 if daash_brands["orders_lw"].sum() > 0 else 0
        chg_o = chg(br.orders_tw, br.orders_lw)
        chg_s = (f"▲ {abs(chg_o):.0f}%" if (chg_o or 0)>0 else f"▼ {abs(chg_o or 0):.0f}%")
        chg_c = GREEN if (chg_o or 0) > 0 else RED
        rows5.append((
            str(rk),
            str(br.brand)[:36],
            f"{int(br.orders_tw):,}",
            f"{share_tw:.1f}%",
            f"{int(br.orders_lw):,}",
            f"{share_lw:.1f}%",
            (chg_s, chg_c),
        ))
    table_block(sl,
                ["#", "Brand / Restaurant", "TW Orders", "TW Share", "LW Orders", "LW Share", "Change"],
                rows5, x=0.28, y=2.52, col_widths=[0.38, 4.05, 1.62, 1.42, 1.62, 1.42, 2.26],
                row_h=0.37, header_color=DAASH_C)

    top = daash_brands.iloc[0]
    top_share = top.orders_tw / total_tw_ord * 100
    top5_share = daash_brands["orders_tw"].sum() / total_tw_ord * 100
    ins5 = (f"'{str(top.brand)[:28]}' leads with {int(top.orders_tw):,} orders "
            f"({top_share:.1f}% share). "
            f"Top brands shown = {top5_share:.1f}% of delivered volume. "
            f"Brands not ordering this week ({r_churned}): follow up — are they on holiday, "
            f"having supply issues, or moving to a competitor?")
    tbl_btm = 2.52 + 0.37 + len(rows5)*0.37 + 0.12
    insight_box(sl, ins5, 0.28, tbl_btm, 12.77, 0.7)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DAASH: ORDER ISSUES
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, LIGHT_BG)
page_header(sl, "DAASH  —  ORDER ISSUES",
            f"Week: {TW_LABEL}   ·   Rejected + voided orders by brand", DAASH_C)

IW, IH, IG = 4.1, 2.05, 0.19
at_r, ac_r = arrow(d_tw_rej,  d_lw_rej,  good="down")
at_v, ac_v = arrow(d_tw_void, d_lw_void, good="down")
at_i, ac_i = arrow(d_tw_irate, d_lw_irate, good="down")

kpi_card(sl, 0.28,            1.22, IW, IH, "Rejected Orders",
         f"{int(d_tw_rej):,}", at_r, ac_r, val_size=38, accent=RED)
kpi_card(sl, 0.28+IW+IG,     1.22, IW, IH, "Voided Orders",
         f"{int(d_tw_void):,}", at_v, ac_v, val_size=38, accent=AMBER)
kpi_card(sl, 0.28+2*(IW+IG), 1.22, IW, IH, "Issue Rate",
         f"{d_tw_irate:.2f}%", at_i, ac_i, val_size=38, accent=MID_GRAY)

# Issue rate banner
ir_bg = RGBColor(0xFE,0xF2,0xF2) if d_tw_irate > 1 else RGBColor(0xF0,0xFD,0xF4)
ir_c  = RED if d_tw_irate > 1 else GREEN
rect(sl, 0.28, 3.43, 12.77, 0.55, fill=ir_bg, border=ir_c, border_pt=1.0)
txt(sl, (f"Issue Rate  {d_tw_irate:.2f}%   ·   "
         f"{int(d_tw_iss):,} of {int(d_tw_total):,} orders affected   ·   "
         f"Industry benchmark: 1–3%"),
    0.5, 3.55, 12.3, 0.35, size=12, bold=True, color=ir_c)

# Issues by brand table
if not daash_issues_by_brand.empty:
    rows6 = []
    for _, r in daash_issues_by_brand.iterrows():
        ir = float(r["issue_rate"])
        ic = RED if ir > 2 else (AMBER if ir > 1 else GREEN)
        rows6.append((
            str(r["brand"])[:32],
            f"{int(r['total_orders']):,}",
            f"{int(r['rejected']):,}",
            f"{int(r['voided']):,}",
            (f"{ir:.1f}%", ic),
        ))
    table_block(sl, ["Brand", "Total Orders", "Rejected", "Voided", "Issue Rate"],
                rows6, x=0.28, y=4.1, col_widths=[4.8, 2.0, 2.0, 1.97, 2.0],
                row_h=0.35, header_color=DAASH_C)
else:
    ir_chg = chg(d_tw_irate, d_lw_irate)
    trend  = ("improving" if (ir_chg or 0) < 0 else "worsening")
    ins6   = (f"Issue rate {d_tw_irate:.2f}% — {trend} WoW. "
              f"{int(d_tw_rej):,} rejected, {int(d_tw_void):,} voided out of "
              f"{int(d_tw_total):,} total orders. "
              f"Pull rejection logs by brand to identify kitchen capacity or acceptance delay issues.")
    insight_box(sl, ins6, 0.28, 4.1, 12.77, 0.7)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — GOSOURCE: WEEK AT A GLANCE
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, LIGHT_BG)
page_header(sl, "GOSOURCE  —  WEEK AT A GLANCE",
            f"Data: {GS_TW_LABEL}{GS_NOTE}   ·   vs Prior: {GS_LW_LABEL}", GS_C)

KW7, KH7, KG7 = 3.03, 2.1, 0.12
for i, (lbl, val, p1, p2, g) in enumerate([
    ("Total Orders",     f"{int(g_tw_ord):,}",  g_tw_ord,  g_lw_ord,  "up"),
    ("Avg Order Value",  naira(g_tw_aov),        g_tw_aov,  g_lw_aov,  "up"),
    ("Active Customers", f"{int(g_tw_cust):,}", g_tw_cust, g_lw_cust, "up"),
    ("Credit Orders",    f"{g_tw_credit_pct:.0f}%", g_tw_credit_pct, None, "up"),
]):
    at, ac = arrow(p1, p2, good=g) if p2 else ("— B2B payment mix", SLATE)
    kpi_card(sl, 0.28 + i*(KW7+KG7), 1.22, KW7, KH7, lbl, val, at, ac,
             val_size=34, accent=GS_C)

# Payment mix visual
pm_y = 3.48
rect(sl, 0.28, pm_y, 12.77, 0.28, fill=GS_C)
txt(sl, "PAYMENT MIX", 0.45, pm_y+0.04, 5, 0.2, size=8.5, bold=True, color=WHITE)

bar_total_w = 12.77
credit_w = bar_total_w * (g_tw_credit_pct / 100)
cash_w    = bar_total_w * (g_tw_cash_pct / 100)
rect(sl, 0.28, pm_y+0.3, credit_w, 0.55, fill=GS_C)
rect(sl, 0.28+credit_w, pm_y+0.3, cash_w, 0.55, fill=TEAL)
txt(sl, f"Credit  {g_tw_credit_pct:.0f}%  ({int(g_tw_credit):,} orders)",
    0.45, pm_y+0.36, 6, 0.38, size=9.5, bold=True, color=WHITE)
txt(sl, f"Cash / Other  {g_tw_cash_pct:.0f}%",
    0.28 + credit_w + 0.12, pm_y+0.36, 5, 0.38, size=9.5, bold=True, color=DARK)

# GoSource pipeline (current month status)
pl_y = 4.55
rect(sl, 0.28, pl_y, 12.77, 0.28, fill=DARK)
txt(sl, "ORDER PIPELINE  —  LAST 30 DAYS", 0.45, pl_y+0.04, 12, 0.2,
    size=8.5, bold=True, color=WHITE)

pl_colors = {"Delivered": GREEN, "Processing": AMBER, "Pending": BLUE,
             "Cancelled": RED,  "Rejected": RED}
if not gs_pipeline.empty:
    PW = 12.77 / len(gs_pipeline)
    for pi, pr in enumerate(gs_pipeline.itertuples()):
        px = 0.28 + pi * PW
        pc = pl_colors.get(str(pr.status), MID_GRAY)
        alt2 = WHITE if pi % 2 == 0 else OFF_WHITE
        rect(sl, px, pl_y+0.3, PW, 1.42, fill=alt2, border=BORDER, border_pt=0.5)
        rect(sl, px, pl_y+0.3, PW, 0.055, fill=pc)
        txt(sl, str(pr.status), px+0.12, pl_y+0.46, PW-0.18, 0.28,
            size=8.5, bold=True, color=pc)
        txt(sl, f"{int(pr.orders):,}", px+0.12, pl_y+0.73, PW-0.18, 0.62,
            size=22, bold=True, color=DARK, wrap=False)
        txt(sl, "orders", px+0.12, pl_y+1.32, PW-0.18, 0.28, size=8, color=SLATE)

g_rev_chg = chg(g_tw_rev, g_lw_rev)
ins7 = (f"{'▲' if (g_rev_chg or 0)>0 else '▼'} {abs(g_rev_chg or 0):.1f}% WoW: "
        f"{int(g_tw_ord):,} delivered orders from {int(g_tw_cust):,} customers "
        f"at {naira(g_tw_aov)} avg. "
        f"Credit mix {g_tw_credit_pct:.0f}% — monitor AR aging closely for credit customers. "
        f"Pipeline shows active orders in system.")
insight_box(sl, ins7, 0.28, 6.08, 12.77, 0.67)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — GOSOURCE: DAY PATTERNS
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, LIGHT_BG)
page_header(sl, "GOSOURCE  —  SALES BY DAY OF WEEK",
            f"Week: {GS_TW_LABEL}   ·   Delivered orders only", GS_C)

if not gs_dow.empty:
    max_gs = int(gs_dow["orders"].max()) or 1
    BAR_TOP = 1.3; BAR_H = 3.8; BAR_MX = BAR_H - 0.7

    col_cnt = len(gs_dow)
    cw8 = 12.77 / col_cnt

    for idx, row in enumerate(gs_dow.itertuples()):
        cx8 = 0.28 + idx * cw8
        o8  = int(row.orders)
        h8  = (o8 / max_gs) * BAR_MX

        alt = RGBColor(0xF3,0xF4,0xF6) if idx % 2 == 0 else WHITE
        rect(sl, cx8, BAR_TOP, cw8, BAR_H + 0.9, fill=alt, border=BORDER, border_pt=0.3)
        txt(sl, str(row.day_name), cx8, BAR_TOP+0.05, cw8, 0.28,
            size=9, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
        bar_y = BAR_TOP + BAR_H - h8 - 0.35
        rect(sl, cx8 + cw8*0.2, bar_y, cw8 * 0.6, h8, fill=GS_C)
        txt(sl, f"{o8}", cx8, BAR_TOP + BAR_H - 0.28, cw8, 0.26,
            size=9, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    if not gs_dow.empty:
        pk = gs_dow.loc[gs_dow["orders"].idxmax()]
        sl8 = gs_dow.loc[gs_dow["orders"].idxmin()]
        ins8 = (f"Peak procurement day: {pk['day_name']} ({int(pk['orders']):,} orders). "
                f"Slowest: {sl8['day_name']} ({int(sl8['orders']):,} orders). "
                f"B2B buying tends to cluster mid-week — use this for route planning and "
                f"supplier dispatch scheduling to ensure on-time delivery for top accounts.")
        insight_box(sl, ins8, 0.28, 5.52, 12.77, 0.72)
else:
    txt(sl, "No day-of-week data available for this period.", 0.28, 2.0, 12.77, 0.5,
        size=12, color=SLATE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — GOSOURCE: AR AGING
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, LIGHT_BG)
page_header(sl, "GOSOURCE  —  ACCOUNTS RECEIVABLE AGING",
            f"Outstanding as of {today.strftime('%d %B %Y')}   ·   Percentage view", GS_C)

bucket_style = {
    "0-30 days":  (GREEN,  RGBColor(0xF0,0xFD,0xF4), RGBColor(0xDC,0xFC,0xE7)),
    "31-60 days": (AMBER,  RGBColor(0xFF,0xFB,0xEB), RGBColor(0xFE,0xF3,0xC7)),
    "61-90 days": (ORANGE, RGBColor(0xFF,0xF7,0xED), RGBColor(0xFE,0xE4,0xC4)),
    "90+ days":   (RED,    RGBColor(0xFE,0xF2,0xF2), RGBColor(0xFE,0xCA,0xCA)),
}

if not gs_ar.empty:
    n   = len(gs_ar)
    bw  = 12.77 / n
    for i, (_, row) in enumerate(gs_ar.iterrows()):
        bkt   = row["ar_aging_bucket"]
        amt   = float(row["amount"])
        share = amt / ar_total * 100 if ar_total > 0 else 0
        tc, bb, bb2 = bucket_style.get(bkt, (MID_GRAY, LIGHT_BG, LIGHT_BG))
        tx = 0.28 + i * bw
        bww = bw - 0.1

        rect(sl, tx, 1.22, bww, 3.6, fill=bb, border=tc, border_pt=1.2)
        rect(sl, tx, 1.22, bww, 0.1, fill=tc)
        txt(sl, bkt.upper(), tx+0.18, 1.42, bww-0.3, 0.3, size=10, bold=True, color=tc)

        # Big % share
        txt(sl, f"{share:.1f}%", tx+0.18, 1.78, bww-0.3, 1.0,
            size=44, bold=True, color=tc, wrap=False)
        txt(sl, "of total AR", tx+0.18, 2.74, bww-0.3, 0.28, size=9, color=MID_GRAY)

        # Amount + count
        txt(sl, naira(amt), tx+0.18, 3.08, bww-0.3, 0.55, size=16, bold=True, color=DARK)
        txt(sl, f"{int(row['invoices']):,} invoices",
            tx+0.18, 3.58, bww-0.3, 0.28, size=9, color=SLATE)

        # Mini progress bar
        bar_ww = (bww - 0.36) * share / 100
        rect(sl, tx+0.18, 3.92, bww-0.36, 0.14, fill=BORDER)
        if bar_ww > 0:
            rect(sl, tx+0.18, 3.92, bar_ww, 0.14, fill=tc)

    # Total banner
    rect(sl, 0.28, 4.95, 12.77, 0.52, fill=DARK)
    txt(sl, f"TOTAL OUTSTANDING  ·  {naira(ar_total)}",
        0.5, 5.03, 12.3, 0.38, size=13, bold=True, color=WHITE)

    ar_ins = (f"{naira(ar_90)} ({ar_90pct:.1f}%) overdue 90+ days — "
              f"immediate collection action required. Send final demand letters and escalate accounts. "
              f"Consider withholding future credit orders until outstanding balances are cleared."
              if ar_90pct > 0 else
              f"AR aging healthy — all {naira(ar_total)} outstanding is within normal collection windows. "
              f"Continue weekly AR reviews and enforce credit terms consistently.")
    ar_bg9 = RGBColor(0xFE,0xF2,0xF2) if ar_90pct > 0 else RGBColor(0xF0,0xFD,0xF4)
    insight_box(sl, ar_ins, 0.28, 5.6, 12.77, 0.68,
                bg_col=ar_bg9, border_col=RED if ar_90pct > 0 else GREEN,
                text_col=RGBColor(0x99,0x1B,0x1B) if ar_90pct > 0 else RGBColor(0x14,0x53,0x2D))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — NOTES & DISCUSSION
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, DARK)
rect(sl, 0, 0, 0.55, 7.5, fill=DAASH_C)
rect(sl, 0.55, 0, 0.18, 7.5, fill=GS_C)

txt(sl, "Notes & Discussion", 1.0, 1.4, 11.5, 0.95, size=30, bold=True, color=WHITE)
txt(sl, "Action items  ·  Key decisions  ·  Observations",
    1.0, 2.3, 11.5, 0.38, size=11, color=RGBColor(0x6B,0x72,0x80))

if GS_TW_DELIVERED == 0:
    txt(sl, f"Note: GoSource current week has 0 delivered orders — slides show {GS_TW_LABEL}",
        1.0, 2.82, 11.5, 0.38, size=9.5, color=RGBColor(0xF8,0x71,0x71))

for i in range(1, 5):
    rect(sl, 1.0, 3.4 + i*0.68, 11.5, 0.56,
         fill=RGBColor(0x1F,0x29,0x3B), border=RGBColor(0x37,0x41,0x51), border_pt=0.5)
    txt(sl, f"{i}.", 1.15, 3.46+i*0.68, 0.42, 0.42, size=11, color=SLATE)

txt(sl, f"IPC Group  ·  Week of {TW_LABEL}  ·  Confidential",
    1.0, 7.1, 11.5, 0.3, size=8.5, color=RGBColor(0x37,0x41,0x51))


# ── Save ───────────────────────────────────────────────────────────────────────
filename = f"IPC_Weekly_Report_{today.strftime('%Y-%m-%d')}.pptx"
output   = os.path.join(os.path.expanduser("~/Desktop"), filename)
prs.save(output)
print(f"\n✅ Saved: {output}")
print(f"   11 slides: Title · Exec Summary · 4× DAASH · 3× GoSource · Notes")
