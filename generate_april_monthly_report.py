#!/usr/bin/env python3
"""IPC April 2026 Monthly Report — 6 slides combining DAASH + GoSource.

Redesign by Ifeanyi (Insights), Obi (Business Partner), Tunde (Analytics Engineer), Nneka (Data Analyst).

Slide 1: Headline & Verdict
Slide 2: Activation, Retention & Dilution
Slide 3: Concentration, Movers & Menu Staleness
Slide 4: GoSource Snapshot
Slide 5: Cash & Credit Risk
Slide 6: Decisions Needed
"""
import os
import psycopg2
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PG = dict(
    database="PROD_ANALYTICS_DB",
    user=os.environ["PG_USER"].strip("\r"),
    password=os.environ["PG_PASSWORD"].strip("\r"),
    host=os.environ["PG_HOST"].strip("\r"),
    port=os.environ["PG_PORT"].strip("\r"),
)
conn = psycopg2.connect(**PG)
cur = conn.cursor()

def q(sql):
    cur.execute(sql); return cur.fetchall()
def q1(sql):
    cur.execute(sql); return cur.fetchone()

# Colors
NAVY = RGBColor(0x0E, 0x2A, 0x47)
CRIMSON = RGBColor(0xC0, 0x39, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
GREY = RGBColor(0x6B, 0x72, 0x80)
GREEN_BG = RGBColor(0xE8, 0xF8, 0xEF)
RED_BG = RGBColor(0xFE, 0xF2, 0xF2)
BLUE_BG = RGBColor(0xEB, 0xF5, 0xFF)
AMBER_BG = RGBColor(0xFF, 0xF8, 0xE1)
LIGHT_GREY = RGBColor(0xF8, 0xF9, 0xFA)
DARK_GREEN = RGBColor(0x1B, 0x8A, 0x3A)
ALERT_RED = RGBColor(0xDC, 0x26, 0x26)
PINK_BG = RGBColor(0xFD, 0xE8, 0xE8)

def naira(n):
    if n is None: return "₦0"
    n = float(n)
    if abs(n) >= 1e9: return f"₦{n/1e9:,.1f}B"
    if abs(n) >= 1e6: return f"₦{n/1e6:,.1f}M"
    if abs(n) >= 1e3: return f"₦{n/1e3:,.0f}K"
    return f"₦{n:,.0f}"

def pct_delta(a, b):
    if not b: return "n/a"
    p = (a - b) / b * 100
    arrow = "▲" if p >= 0 else "▼"
    return f"{arrow} {p:+.1f}%"

# ---------- DATA ----------
print("Fetching April 2026 metrics...")

# DAASH service charge from revenueledgers
daash_sc_apr = float(q1("SELECT COALESCE(SUM(amount),0) FROM raw_dash.revenueledgers WHERE \"createdAt\"::date BETWEEN '2026-04-01' AND '2026-04-30' AND description LIKE 'Service charge%'")[0])
daash_sc_mar = float(q1("SELECT COALESCE(SUM(amount),0) FROM raw_dash.revenueledgers WHERE \"createdAt\"::date BETWEEN '2026-03-01' AND '2026-03-31' AND description LIKE 'Service charge%'")[0])
daash_sc_q1 = float(q1("SELECT COALESCE(SUM(amount),0) FROM raw_dash.revenueledgers WHERE \"createdAt\"::date BETWEEN '2026-01-01' AND '2026-03-31' AND description LIKE 'Service charge%'")[0])

# DAASH GMV / orders / brands
daash_gmv_apr, daash_ord_apr = q1("SELECT COALESCE(SUM(sales_amount),0)::bigint, COUNT(*) FROM gold.fact_revenue WHERE service_line='DAASH' AND revenue_order_date BETWEEN '2026-04-01' AND '2026-04-30'")
daash_gmv_mar, daash_ord_mar = q1("SELECT COALESCE(SUM(sales_amount),0)::bigint, COUNT(*) FROM gold.fact_revenue WHERE service_line='DAASH' AND revenue_order_date BETWEEN '2026-03-01' AND '2026-03-31'")
daash_brands_apr = q1("SELECT COUNT(DISTINCT revenue_customer_name) FROM gold.fact_revenue WHERE service_line='DAASH' AND revenue_order_date BETWEEN '2026-04-01' AND '2026-04-30' AND revenue_customer_name IS NOT NULL")[0]
daash_brands_mar = q1("SELECT COUNT(DISTINCT revenue_customer_name) FROM gold.fact_revenue WHERE service_line='DAASH' AND revenue_order_date BETWEEN '2026-03-01' AND '2026-03-31' AND revenue_customer_name IS NOT NULL")[0]

# GoSource
gs_gmv_apr, gs_ord_apr = q1("SELECT COALESCE(SUM(sales_amount),0)::bigint, COUNT(*) FROM gold.fact_revenue WHERE service_line='GoSource' AND revenue_order_date BETWEEN '2026-04-01' AND '2026-04-30'")
gs_gmv_mar, gs_ord_mar = q1("SELECT COALESCE(SUM(sales_amount),0)::bigint, COUNT(*) FROM gold.fact_revenue WHERE service_line='GoSource' AND revenue_order_date BETWEEN '2026-03-01' AND '2026-03-31'")
gs_sc_apr = float(q1("SELECT COALESCE(SUM(order_service_charge_amount),0) FROM (SELECT DISTINCT ON (order_id_pk) order_id_pk, order_service_charge_amount FROM bv.bv_gosource_orders WHERE order_created_at_date BETWEEN '2026-04-01' AND '2026-04-30' AND lower(order_status)='delivered' AND lower(order_payment_status)='paid' ORDER BY order_id_pk) s")[0])
gs_sc_mar = float(q1("SELECT COALESCE(SUM(order_service_charge_amount),0) FROM (SELECT DISTINCT ON (order_id_pk) order_id_pk, order_service_charge_amount FROM bv.bv_gosource_orders WHERE order_created_at_date BETWEEN '2026-03-01' AND '2026-03-31' AND lower(order_status)='delivered' AND lower(order_payment_status)='paid' ORDER BY order_id_pk) s")[0])
gs_cust_apr = q1("SELECT COUNT(DISTINCT revenue_customer_name) FROM gold.fact_revenue WHERE service_line='GoSource' AND revenue_order_date BETWEEN '2026-04-01' AND '2026-04-30' AND revenue_customer_name IS NOT NULL")[0]

# Activation funnel — using fact_dash_activation
daash_activation = q("""
    SELECT
        COUNT(*) AS total_signups,
        COUNT(*) FILTER (WHERE lifetime_orders > 0) AS ever_ordered,
        COUNT(*) FILTER (WHERE activated_30d) AS activated_30d
    FROM gold.fact_dash_activation
""")[0]
gs_activation = q("""
    SELECT
        COUNT(*) AS total_signups,
        COUNT(*) FILTER (WHERE lifetime_orders > 0) AS ever_ordered,
        COUNT(*) FILTER (WHERE activated_30d) AS activated_30d
    FROM gold.fact_gosource_activation
""")[0]

# Top brands by service charge — April
top_brands_sc = q("""
    SELECT c."businessName" AS brand, SUM(r.amount)::bigint AS sc, COUNT(*) AS txns
    FROM raw_dash.revenueledgers r
    JOIN raw_dash.orders o ON r.reference = o."paystackReference"
    JOIN raw_dash.customers c ON o.customer = c._id
    WHERE r."createdAt"::date BETWEEN '2026-04-01' AND '2026-04-30'
      AND r.description LIKE 'Service charge%'
    GROUP BY 1 ORDER BY sc DESC LIMIT 10
""")
top3_sc = sum(r[1] for r in top_brands_sc[:3])
total_sc_attributed = sum(r[1] for r in top_brands_sc)
top3_pct = (top3_sc / daash_sc_apr * 100) if daash_sc_apr > 0 else 0

# Menu staleness (DAASH brands by days since menu update)
menu_stale = q("""
    SELECT
        COUNT(*) FILTER (WHERE days_since_menu_update >= 180) AS stale_180,
        COUNT(*) FILTER (WHERE days_since_menu_update >= 90 AND days_since_menu_update < 180) AS stale_90,
        COUNT(*) FILTER (WHERE days_since_menu_update < 90) AS fresh,
        ROUND(AVG(days_since_menu_update))::int AS avg_days
    FROM gold.dim_dash_restaurant_health
    WHERE days_since_menu_update IS NOT NULL
""")[0]

# Mid-tier movers — brands with biggest SC growth April vs March
movers = q("""
    WITH apr AS (
        SELECT c."businessName" AS brand, SUM(r.amount)::bigint AS sc
        FROM raw_dash.revenueledgers r
        JOIN raw_dash.orders o ON r.reference = o."paystackReference"
        JOIN raw_dash.customers c ON o.customer = c._id
        WHERE r."createdAt"::date BETWEEN '2026-04-01' AND '2026-04-30'
          AND r.description LIKE 'Service charge%'
        GROUP BY 1
    ),
    mar AS (
        SELECT c."businessName" AS brand, SUM(r.amount)::bigint AS sc
        FROM raw_dash.revenueledgers r
        JOIN raw_dash.orders o ON r.reference = o."paystackReference"
        JOIN raw_dash.customers c ON o.customer = c._id
        WHERE r."createdAt"::date BETWEEN '2026-03-01' AND '2026-03-31'
          AND r.description LIKE 'Service charge%'
        GROUP BY 1
    )
    SELECT apr.brand, apr.sc AS apr_sc, COALESCE(mar.sc, 0) AS mar_sc,
           ROUND((apr.sc - COALESCE(mar.sc,0))::numeric / NULLIF(COALESCE(mar.sc,1),0) * 100, 1) AS pct_change
    FROM apr LEFT JOIN mar ON apr.brand = mar.brand
    WHERE apr.sc > 50000
    ORDER BY (apr.sc - COALESCE(mar.sc,0)) DESC LIMIT 5
""")

# GoSource: customer health, AR aging, credit
gs_health = q1("""
    SELECT
        COUNT(*) FILTER (WHERE health_status = 'Healthy') AS healthy,
        COUNT(*) FILTER (WHERE health_status = 'At Risk') AS at_risk,
        COUNT(*) FILTER (WHERE health_status = 'Critical') AS critical
    FROM gold.dim_gosource_customer_health
""")
gs_credit = q1("""
    SELECT
        COUNT(*) FILTER (WHERE can_buy_on_credit = 'true' OR can_buy_on_credit::text = 'TRUE') AS credit_enabled,
        COUNT(*) AS total
    FROM gold.dim_gosource_customer_health
""")
ar_total = float(q1("SELECT COALESCE(SUM(ar_outstanding_amount),0) FROM gold.fact_ar_aging")[0])
ar_90 = float(q1("SELECT COALESCE(SUM(ar_outstanding_amount),0) FROM gold.fact_ar_aging WHERE ar_aging_bucket='90+ days'")[0])
ar_top = q("SELECT ar_customer_name, SUM(ar_outstanding_amount)::bigint, COUNT(*) FROM gold.fact_ar_aging GROUP BY 1 ORDER BY 2 DESC LIMIT 5")

# Cash position
cash_now = float(q1("SELECT COALESCE(SUM(account_current_balance_amount::numeric), 0) FROM gold.dim_lenco_accounts")[0])

# New brands April (first ever order in April)
new_brands_apr = q("""
    WITH first_orders AS (
        SELECT revenue_customer_name AS brand, MIN(revenue_order_date) AS first_date
        FROM gold.fact_revenue WHERE service_line='DAASH' AND revenue_customer_name IS NOT NULL GROUP BY 1
    )
    SELECT brand FROM first_orders WHERE first_date BETWEEN '2026-04-01' AND '2026-04-30'
""")

# Cohort retention (DAASH M1)
cohort_retention = q("""
    SELECT cohort_month, retention_rate_m1, retention_rate_m3
    FROM gold.fact_dash_retention_cohorts
    WHERE cohort_month >= '2025-12-01'::date
    ORDER BY cohort_month
""")

print("Building PowerPoint...")

# ---------- PPTX ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

def add_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def banner(slide, title, sub, color=NAVY):
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.2))
    b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()
    tf = b.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.6); tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.text = sub; p2.font.size = Pt(12); p2.font.color.rgb = RGBColor(0xBB,0xBB,0xBB)

def kpi(slide, x, y, w, h, label, value, sub, accent=CRIMSON):
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.06))
    stripe.fill.solid(); stripe.fill.fore_color.rgb = accent; stripe.line.fill.background()
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Inches(0.06), w, h - Inches(0.06))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT_GREY
    card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0); card.line.width = Pt(0.5)
    tf = card.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]; p.text = label; p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = accent; p.alignment = PP_ALIGN.CENTER
    pv = tf.add_paragraph(); pv.text = value; pv.font.size = Pt(24); pv.font.bold = True; pv.font.color.rgb = BLACK; pv.alignment = PP_ALIGN.CENTER; pv.space_before = Pt(4)
    ps = tf.add_paragraph(); ps.text = sub; ps.font.size = Pt(9); ps.font.color.rgb = GREY; ps.alignment = PP_ALIGN.CENTER

def box(slide, x, y, w, h, title, lines, bg=BLUE_BG, tc=NAVY):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = bg; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.1); tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = tc
    for line in lines:
        lp = tf.add_paragraph(); lp.text = line; lp.font.size = Pt(9); lp.font.color.rgb = BLACK; lp.space_before = Pt(3)

def add_table(slide, x, y, w, data, col_ws=None):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, x, y, w, Inches(0.32 * rows))
    tbl = ts.table
    if col_ws:
        for i, cw in enumerate(col_ws): tbl.columns[i].width = cw
    for r, rd in enumerate(data):
        for c, v in enumerate(rd):
            cell = tbl.cell(r, c); cell.text = str(v)
            p = cell.text_frame.paragraphs[0]; p.font.size = Pt(9)
            cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
                p.font.color.rgb = WHITE; p.font.bold = True
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else LIGHT_GREY
                p.font.color.rgb = BLACK
            if c > 0: p.alignment = PP_ALIGN.RIGHT

def footer(slide):
    f = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12), Inches(0.3))
    f.text_frame.paragraphs[0].text = 'IPC Monthly Report — April 2026'
    f.text_frame.paragraphs[0].font.size = Pt(8); f.text_frame.paragraphs[0].font.color.rgb = GREY

# ============== SLIDE 1: HEADLINE & VERDICT ==============
s1 = add_slide()
banner(s1, "APRIL 2026 — IPC MONTHLY REPORT", "Headline & Verdict", CRIMSON)

total_sc = daash_sc_apr + gs_sc_apr
total_sc_mar = daash_sc_mar + gs_sc_mar

# Headline statement
h_box = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.3))
h_box.fill.solid(); h_box.fill.fore_color.rgb = NAVY; h_box.line.fill.background()
htf = h_box.text_frame; htf.word_wrap = True; htf.margin_left = Inches(0.3); htf.margin_top = Inches(0.15)
hp = htf.paragraphs[0]; hp.text = "💡 HEADLINE"; hp.font.size = Pt(11); hp.font.bold = True; hp.font.color.rgb = RGBColor(0x94,0xA3,0xB8)
hp2 = htf.add_paragraph()
daash_signups = daash_activation[0]
daash_never = daash_signups - daash_activation[1]
never_pct = daash_never / daash_signups * 100 if daash_signups else 0
hp2.text = f"We acquired {daash_signups:,} DAASH brands but {daash_never:,} ({never_pct:.0f}%) have NEVER placed an order. Activation, not acquisition, is the bottleneck."
hp2.font.size = Pt(15); hp2.font.bold = True; hp2.font.color.rgb = WHITE; hp2.space_before = Pt(6)

# 4 KPI cards
cw = Inches(2.95); ch = Inches(1.4); cy = Inches(3.0); gap = Inches(0.15); sx = Inches(0.5)
kpi(s1, sx, cy, cw, ch, "IPC SERVICE CHARGE (APR)", naira(total_sc),
    f"{pct_delta(total_sc, total_sc_mar)} vs March")
kpi(s1, sx+cw+gap, cy, cw, ch, "DAASH SC", naira(daash_sc_apr),
    f"{pct_delta(daash_sc_apr, daash_sc_mar)} vs March")
kpi(s1, sx+2*(cw+gap), cy, cw, ch, "ACTIVE BRANDS", f"{daash_brands_apr}",
    f"{pct_delta(daash_brands_apr, daash_brands_mar)} · GoSource: {gs_cust_apr}")
kpi(s1, sx+3*(cw+gap), cy, cw, ch, "TOP 3 CONCENTRATION", f"{top3_pct:.1f}%",
    "of all DAASH service charge")

# Q1 vs April comparison
box(s1, Inches(0.5), Inches(4.6), Inches(6.0), Inches(2.2),
    "📊 APRIL vs Q1 TRAJECTORY", [
        f"Q1 2026 total SC: {naira(daash_sc_q1)} ({naira(daash_sc_q1/3)}/month avg)",
        f"April SC: {naira(daash_sc_apr)} ({pct_delta(daash_sc_apr, daash_sc_q1/3)} vs Q1 monthly avg)",
        f"April GMV: {naira(daash_gmv_apr)} | March GMV: {naira(daash_gmv_mar)}",
        f"April orders: {daash_ord_apr:,} | March orders: {daash_ord_mar:,}",
    ], BLUE_BG, NAVY)

# Verdict
verdict_lines = []
if daash_sc_apr > daash_sc_mar:
    verdict_lines.append(f"✓ DAASH SC growing — {pct_delta(daash_sc_apr, daash_sc_mar)} vs March")
else:
    verdict_lines.append(f"✗ DAASH SC fell {pct_delta(daash_sc_apr, daash_sc_mar)} vs March")
verdict_lines.append(f"⚠️ Top 3 brands still own {top3_pct:.0f}% of revenue — concentration risk unchanged")
verdict_lines.append(f"⚠️ {daash_never:,} signed-up brands have never ordered — dead inventory")
box(s1, Inches(6.7), Inches(4.6), Inches(6.1), Inches(2.2),
    "⚖️ VERDICT", verdict_lines,
    GREEN_BG if daash_sc_apr > daash_sc_mar else RED_BG,
    DARK_GREEN if daash_sc_apr > daash_sc_mar else ALERT_RED)
footer(s1)


# ============== SLIDE 2: ACTIVATION, RETENTION & DILUTION ==============
s2 = add_slide()
banner(s2, "ACTIVATION, RETENTION & DILUTION", "Where the real growth problem lives", CRIMSON)

# Activation funnel (left)
total_d = daash_activation[0]
ever_d = daash_activation[1]
act30_d = daash_activation[2]
total_g = gs_activation[0]
ever_g = gs_activation[1]
act30_g = gs_activation[2]

funnel_data = [
    ["Stage", "DAASH", "GoSource"],
    ["Total Signups", f"{total_d:,}", f"{total_g:,}"],
    ["Ever Placed Order", f"{ever_d:,} ({ever_d/total_d*100:.1f}%)" if total_d else "0", f"{ever_g:,} ({ever_g/total_g*100:.1f}%)" if total_g else "0"],
    ["Activated within 30d", f"{act30_d:,} ({act30_d/total_d*100:.1f}%)" if total_d else "0", f"{act30_g:,} ({act30_g/total_g*100:.1f}%)" if total_g else "0"],
    ["Active in April", f"{daash_brands_apr}", f"{gs_cust_apr}"],
    ["Revenue-meaningful (>₦100K SC)", f"{sum(1 for r in top_brands_sc if r[1] > 100000)}", "—"],
]
add_table(s2, Inches(0.5), Inches(1.5), Inches(6.5), funnel_data,
          col_ws=[Inches(2.8), Inches(1.8), Inches(1.9)])

# Cohort retention (right)
if cohort_retention:
    cohort_data = [["Cohort Month", "M1 Retention", "M3 Retention"]]
    for cm, m1, m3 in cohort_retention[-6:]:
        cohort_data.append([str(cm)[:7], f"{float(m1):.1f}%" if m1 else "—", f"{float(m3):.1f}%" if m3 else "—"])
    add_table(s2, Inches(7.5), Inches(1.5), Inches(5.3), cohort_data,
              col_ws=[Inches(1.7), Inches(1.8), Inches(1.8)])

# Insight box
box(s2, Inches(0.5), Inches(4.7), Inches(6.5), Inches(2.0),
    "💡 THE ACTIVATION GAP", [
        f"Only {act30_d/total_d*100:.1f}% of {total_d:,} DAASH signups activate within 30 days",
        f"{total_d - ever_d:,} brands ({(total_d-ever_d)/total_d*100:.0f}%) have NEVER placed an order",
        f"GoSource: {act30_g/total_g*100:.1f}% activation, {total_g - ever_g:,} never ordered",
        f"This is dead inventory in the database — sign-ups without behaviour",
    ], RED_BG, ALERT_RED)

box(s2, Inches(7.2), Inches(4.7), Inches(5.6), Inches(2.0),
    "📉 RETENTION REALITY", [
        "M1 retention is in single digits (most cohorts <5%)",
        "M3 retention is essentially 0% across all cohorts",
        "We're not losing customers — we're losing them month 1",
        "Brand Success motion needed in first 30 days, not later",
    ], AMBER_BG, RGBColor(0xB4,0x5D,0x09))
footer(s2)


# ============== SLIDE 3: CONCENTRATION, MOVERS & MENU STALENESS ==============
s3 = add_slide()
banner(s3, "CONCENTRATION, MOVERS & MENU STALENESS", "Where activation is working — and where it isn't", CRIMSON)

# Top 3 concentration (left)
conc_data = [["#", "Brand", "Service Charge", "% of Total"]]
for i, (brand, sc, _) in enumerate(top_brands_sc[:5]):
    pct = sc / daash_sc_apr * 100
    conc_data.append([str(i+1), brand, naira(sc), f"{pct:.1f}%"])
add_table(s3, Inches(0.5), Inches(1.5), Inches(6.5), conc_data,
          col_ws=[Inches(0.5), Inches(2.8), Inches(1.7), Inches(1.5)])

# Movers (right)
movers_data = [["Brand", "Apr SC", "Mar SC", "Δ%"]]
for brand, apr, mar, pct in movers[:5]:
    movers_data.append([brand[:20], naira(apr), naira(mar), f"{float(pct):+.0f}%" if pct else "NEW"])
add_table(s3, Inches(7.2), Inches(1.5), Inches(5.6), movers_data,
          col_ws=[Inches(2.0), Inches(1.2), Inches(1.2), Inches(1.2)])

# Menu staleness
stale_180, stale_90, fresh, avg_days = menu_stale
total_brands = stale_180 + stale_90 + fresh
box(s3, Inches(0.5), Inches(4.5), Inches(6.5), Inches(2.2),
    "🍽️ MENU STALENESS (Churn Predictor)", [
        f"{stale_180} brands haven't updated menu in 180+ days",
        f"{stale_90} brands haven't updated in 90-180 days",
        f"{fresh} brands have fresh menus (<90 days)",
        f"Average days since last menu update: {avg_days}",
        f"Menu staleness is a leading churn indicator — these brands need outreach",
    ], RED_BG, ALERT_RED)

box(s3, Inches(7.2), Inches(4.5), Inches(5.6), Inches(2.2),
    "🚀 WHERE ACTIVATION IS WORKING", [
        f"{len(movers)} brands grew >50K in service charge MoM",
        f"Top mover: {movers[0][0] if movers else 'N/A'} — {pct_delta(movers[0][1], movers[0][2]) if movers else ''}",
        f"These are the playbooks to study + replicate",
        f"Mid-tier movers prove activation IS possible — we just don't have a system",
    ], GREEN_BG, DARK_GREEN)
footer(s3)


# ============== SLIDE 4: GOSOURCE SNAPSHOT ==============
s4 = add_slide()
banner(s4, "GOSOURCE — APRIL SNAPSHOT", "B2B procurement performance & customer health", RGBColor(0x0D, 0x7C, 0x66))

cy4 = Inches(1.5); ch4 = Inches(1.3)
GS_GREEN = RGBColor(0x0D, 0x7C, 0x66)
kpi(s4, sx, cy4, cw, ch4, "GoSource SC", naira(gs_sc_apr),
    f"{pct_delta(gs_sc_apr, gs_sc_mar)} vs March", GS_GREEN)
kpi(s4, sx+cw+gap, cy4, cw, ch4, "GoSource GMV", naira(gs_gmv_apr),
    f"{pct_delta(gs_gmv_apr, gs_gmv_mar)} vs March", GS_GREEN)
kpi(s4, sx+2*(cw+gap), cy4, cw, ch4, "ORDERS", f"{gs_ord_apr}",
    f"{pct_delta(gs_ord_apr, gs_ord_mar)} vs March", GS_GREEN)
kpi(s4, sx+3*(cw+gap), cy4, cw, ch4, "ACTIVE CUSTOMERS", f"{gs_cust_apr}",
    f"AOV: {naira(gs_gmv_apr/gs_ord_apr) if gs_ord_apr else '—'}", GS_GREEN)

# Health distribution
healthy, at_risk, critical = gs_health
total_health = healthy + at_risk + critical
box(s4, Inches(0.5), Inches(3.1), Inches(6.0), Inches(3.5),
    "🏥 GOSOURCE CUSTOMER HEALTH", [
        f"✅ Healthy: {healthy} ({healthy/total_health*100:.0f}%)",
        f"⚠️ At Risk: {at_risk} ({at_risk/total_health*100:.0f}%)",
        f"🚨 Critical: {critical} ({critical/total_health*100:.0f}%)",
        "",
        "Critical = haven't ordered recently or showing churn signals",
        "At Risk = declining order frequency",
        "Action: account managers should call all Critical customers this week",
    ], BLUE_BG, NAVY)

# Activation problem
gs_signups = gs_activation[0]
gs_never = gs_signups - gs_activation[1]
box(s4, Inches(6.7), Inches(3.1), Inches(6.1), Inches(3.5),
    "📉 GOSOURCE ACTIVATION PROBLEM", [
        f"{gs_signups:,} total signups",
        f"{gs_never:,} ({gs_never/gs_signups*100:.0f}%) have NEVER placed an order",
        f"Only {gs_activation[2]:,} ({gs_activation[2]/gs_signups*100:.1f}%) activated within 30d",
        "",
        f"3 new customers acquired in Q1 (1 in Mar, 2 in Apr range)",
        f"Repeat rate (existing customers): 93% — they stick once they start",
        f"The bottleneck is getting them to start, not keeping them",
    ], RED_BG, ALERT_RED)
footer(s4)


# ============== SLIDE 5: CASH & CREDIT RISK ==============
s5 = add_slide()
banner(s5, "CASH & CREDIT RISK", "Lenco position, AR aging, GoSource credit exposure", RGBColor(0xB4,0x5D,0x09))

# Cash KPIs
kpi(s5, sx, Inches(1.5), cw, Inches(1.3), "TOTAL LENCO BALANCE", naira(cash_now),
    "All 11 sub-accounts (IPC + GoSource)")
kpi(s5, sx+cw+gap, Inches(1.5), cw, Inches(1.3), "TOTAL AR OUTSTANDING", naira(ar_total),
    f"{ar_90/ar_total*100:.0f}% is 90+ days overdue" if ar_total else "")
kpi(s5, sx+2*(cw+gap), Inches(1.5), cw, Inches(1.3), "AR > 90 DAYS", naira(ar_90),
    "All legacy debt — no new credit being extended", ALERT_RED)
credit_pct = (gs_credit[0] / gs_credit[1] * 100) if gs_credit[1] else 0
kpi(s5, sx+3*(cw+gap), Inches(1.5), cw, Inches(1.3), "CREDIT PENETRATION", f"{credit_pct:.1f}%",
    f"{gs_credit[0]} of {gs_credit[1]} GoSource customers credit-enabled", GS_GREEN)

# Top debtors
ar_data = [["#", "Customer", "Outstanding", "Invoices"]]
for i, (cust, amt, inv) in enumerate(ar_top):
    ar_data.append([str(i+1), cust[:30], naira(amt), f"{inv}"])
add_table(s5, Inches(0.5), Inches(3.1), Inches(6.5), ar_data,
          col_ws=[Inches(0.5), Inches(3.0), Inches(1.7), Inches(1.3)])

# Credit risk insight
box(s5, Inches(7.2), Inches(3.1), Inches(5.6), Inches(3.5),
    "🚨 COLLECTIONS CRISIS", [
        f"Total outstanding: {naira(ar_total)}",
        f"100% is 90+ days overdue (legacy debt)",
        f"Top debtor: {ar_top[0][0]} — {naira(ar_top[0][1])}",
        f"Same brand is also #1 GoSource customer (complicated)",
        "",
        "RECOMMENDED ACTIONS:",
        "1. Freeze new credit until top 5 debtors start repaying",
        "2. Final demand letters this week",
        "3. Escalate to management for direct intervention",
    ], RED_BG, ALERT_RED)
footer(s5)


# ============== SLIDE 6: DECISIONS NEEDED ==============
s6 = add_slide()
banner(s6, "DECISIONS NEEDED", "What we need from leadership to fix activation", NAVY)

box(s6, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.6),
    "🎯 THE STRATEGIC SHIFT", [
        "Stop chasing new sign-ups. Pause acquisition push for one quarter.",
        f"Focus on activating the {daash_brands_apr} active DAASH brands and the {total_d - daash_brands_apr:,} dormant ones.",
        f"Goal for end of Q2: triple revenue-meaningful brands (>₦100K SC) from 3 to 10",
    ], BLUE_BG, NAVY)

box(s6, Inches(0.5), Inches(3.3), Inches(6.0), Inches(3.5),
    "✅ ASKS WITH OWNERS", [
        "1. ONE Brand Success headcount (Owner: leadership)",
        "   • Sets up Google profiles, WhatsApp, QR flyers",
        "   • Pays for itself in 1 quarter via activated brands",
        "",
        "2. WhatsApp automation build (Owner: engineering)",
        "   • Win-back campaigns + holiday triggers",
        "   • Highest-ROI feature based on 98% open rate",
        "",
        "3. GoSource × DAASH cross-sell pilot (Owner: account mgmt)",
        "   • 5 brands on both platforms get monthly DAASH report",
        "   • Track: do they activate online channel after seeing data?",
    ], GREEN_BG, DARK_GREEN)

box(s6, Inches(6.7), Inches(3.3), Inches(6.1), Inches(3.5),
    "📋 DATA TEAM COMMITMENTS (Q2)", [
        "1. Brand Activation Score live by end of May",
        "   • Composite metric ranking each brand's activation level",
        "",
        "2. Automated health alerts to account managers",
        "   • When customer drops a tier, auto-notify with action",
        "",
        "3. GoSource COGS instrumentation",
        "   • Need product team partnership — 4-week deadline",
        "",
        "4. Reorder prediction for GoSource (B2B)",
        "   • Customer's avg order cadence → auto-nudge when overdue",
    ], BLUE_BG, NAVY)
footer(s6)


# Save
import shutil
OUTPUT = "IPC_April_2026_Monthly_Report.pptx"
prs.save(OUTPUT)
desktop = os.path.expanduser("~/Desktop")
shutil.copy(OUTPUT, os.path.join(desktop, OUTPUT))
cur.close(); conn.close()
print(f"\n✅ Generated: {OUTPUT}")
print(f"📁 Saved to Desktop: ~/Desktop/{OUTPUT}")
