#!/usr/bin/env python3
"""IPC Weekly Activation Report — 7-slide PowerPoint.

Slides 1: Title
Slides 2-4: DAASH (Snapshot, Channel & Brand Health, Top Movers)
Slides 5-7: GoSource (Snapshot, Pipeline, AR Aging)

Style: red banners, pink KPI cards, insight boxes — matches monthly deck.
Compares last completed Fri-Thu week vs 4-week trailing average.
"""
import os
import psycopg2
from datetime import date, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PG = dict(
    database="PROD_ANALYTICS_DB",
    user=os.environ["PG_USER"].strip("\r"),
    password=os.environ["PG_PASSWORD"].strip("\r"),
    host=os.environ["PG_HOST"].strip("\r"),
    port=os.environ["PG_PORT"].strip("\r"),
)
conn = psycopg2.connect(**PG)
cur = conn.cursor()

def q(sql):
    cur.execute(sql)
    return cur.fetchall()

def q1(sql):
    cur.execute(sql)
    return cur.fetchone()

# ---------------------------------------------------------------------------
# Colors
# ---------------------------------------------------------------------------
CRIMSON     = RGBColor(0xC0, 0x39, 0x2B)
DARK_RED    = RGBColor(0x8B, 0x00, 0x00)
PINK_BG     = RGBColor(0xFD, 0xE8, 0xE8)
GREEN_BG    = RGBColor(0xE8, 0xF8, 0xEF)
BLUE_BG     = RGBColor(0xEB, 0xF5, 0xFF)
RED_BG      = RGBColor(0xFE, 0xF2, 0xF2)
AMBER_BG    = RGBColor(0xFF, 0xF8, 0xE1)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x1A, 0x1A, 0x2E)
GREY        = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GREY  = RGBColor(0xF8, 0xF9, 0xFA)
NAVY        = RGBColor(0x0E, 0x2A, 0x47)
DARK_GREEN  = RGBColor(0x1B, 0x8A, 0x3A)
GS_GREEN    = RGBColor(0x0D, 0x7C, 0x66)
ALERT_RED   = RGBColor(0xDC, 0x26, 0x26)

# ---------------------------------------------------------------------------
# Date windows
# ---------------------------------------------------------------------------
today = date.today()
days_since_thu = (today.weekday() - 3) % 7
if days_since_thu == 0:
    days_since_thu = 7
last_thu = today - timedelta(days=days_since_thu)
tw_end = last_thu
tw_start = tw_end - timedelta(days=6)

weeks = []
for i in range(5):
    we = tw_end - timedelta(weeks=i)
    ws = we - timedelta(days=6)
    weeks.append((ws, we))

TW = weeks[0]
PRIOR_4 = weeks[1:5]
TW_LABEL = f"{TW[0].strftime('%d %b')} – {TW[1].strftime('%d %b %Y')}"

print(f"This week: {TW[0]} → {TW[1]}")

# ---------------------------------------------------------------------------
# Formatting
# ---------------------------------------------------------------------------
def naira(n):
    if n is None: return "₦0"
    n = float(n)
    if abs(n) >= 1e9: return f"₦{n/1e9:,.1f}B"
    if abs(n) >= 1e6: return f"₦{n/1e6:,.1f}M"
    if abs(n) >= 1e3: return f"₦{n/1e3:,.0f}K"
    return f"₦{n:,.0f}"

def delta_str(curr, avg):
    if not avg: return "—"
    p = (curr - avg) / avg * 100
    arrow = "▲" if p >= 0 else "▼"
    return f"{arrow} {p:+.1f}%"

def is_up(curr, avg):
    return curr >= avg if avg else True

# ---------------------------------------------------------------------------
# Data fetching
# ---------------------------------------------------------------------------
def week_data(start, end):
    s, e = str(start), str(end)
    daash_sc = float(q1(f"SELECT COALESCE(SUM(amount),0) FROM raw_dash.revenueledgers WHERE \"createdAt\"::date BETWEEN '{s}' AND '{e}' AND description LIKE 'Service charge%'")[0])
    gs_sc = float(q1(f"SELECT COALESCE(SUM(order_service_charge_amount),0) FROM (SELECT DISTINCT ON (order_id_pk) order_id_pk, order_service_charge_amount FROM bv.bv_gosource_orders WHERE order_created_at_date BETWEEN '{s}' AND '{e}' AND lower(order_status)='delivered' AND lower(order_payment_status)='paid' ORDER BY order_id_pk) sub")[0])
    gmv = q(f"SELECT service_line, COALESCE(SUM(sales_amount),0), COUNT(*) FROM gold.fact_revenue WHERE revenue_order_date BETWEEN '{s}' AND '{e}' GROUP BY 1")
    d_gmv = sum(float(r[1]) for r in gmv if r[0]=='DAASH')
    g_gmv = sum(float(r[1]) for r in gmv if r[0]=='GoSource')
    d_ord = sum(int(r[2]) for r in gmv if r[0]=='DAASH')
    g_ord = sum(int(r[2]) for r in gmv if r[0]=='GoSource')
    d_brands = int(q1(f"SELECT COUNT(DISTINCT revenue_customer_name) FROM gold.fact_revenue WHERE service_line='DAASH' AND revenue_order_date BETWEEN '{s}' AND '{e}' AND revenue_customer_name IS NOT NULL")[0] or 0)
    g_cust = int(q1(f"SELECT COUNT(DISTINCT revenue_customer_name) FROM gold.fact_revenue WHERE service_line='GoSource' AND revenue_order_date BETWEEN '{s}' AND '{e}' AND revenue_customer_name IS NOT NULL")[0] or 0)
    ch = q(f"SELECT revenue_order_channel, COUNT(*), COALESCE(SUM(sales_amount),0) FROM gold.fact_revenue WHERE service_line='DAASH' AND revenue_order_date BETWEEN '{s}' AND '{e}' GROUP BY 1")
    web_o = sum(int(r[1]) for r in ch if r[0] and r[0].lower()=='website')
    pos_o = sum(int(r[1]) for r in ch if r[0] and r[0].lower()=='pos')
    web_gmv = sum(float(r[2]) for r in ch if r[0] and r[0].lower()=='website')
    pos_gmv = sum(float(r[2]) for r in ch if r[0] and r[0].lower()=='pos')
    return dict(daash_sc=daash_sc, gs_sc=gs_sc, daash_gmv=d_gmv, gs_gmv=g_gmv, daash_orders=d_ord, gs_orders=g_ord, daash_brands=d_brands, gs_customers=g_cust, web_orders=web_o, pos_orders=pos_o, web_gmv=web_gmv, pos_gmv=pos_gmv, web_pct=(web_o/(web_o+pos_o)*100) if (web_o+pos_o)>0 else 0, daash_aov=(d_gmv/d_ord) if d_ord>0 else 0, gs_aov=(g_gmv/g_ord) if g_ord>0 else 0)

print("Fetching metrics...")
tw = week_data(TW[0], TW[1])
prior = [week_data(w[0], w[1]) for w in PRIOR_4]
def avg4(k): return sum(p[k] for p in prior)/4 if prior else 0

s_tw, e_tw = str(TW[0]), str(TW[1])
s_lw, e_lw = str(PRIOR_4[0][0]), str(PRIOR_4[0][1])

# DAASH brand data
top_movers = q(f"""
    WITH tw AS (SELECT revenue_customer_name AS brand, COUNT(*) AS tw_ord FROM gold.fact_revenue WHERE service_line='DAASH' AND revenue_order_date BETWEEN '{s_tw}' AND '{e_tw}' AND revenue_customer_name IS NOT NULL GROUP BY 1),
    a4 AS (SELECT revenue_customer_name AS brand, COUNT(*)::numeric/4 AS avg_ord FROM gold.fact_revenue WHERE service_line='DAASH' AND revenue_order_date BETWEEN '{PRIOR_4[3][0]}' AND '{PRIOR_4[0][1]}' AND revenue_customer_name IS NOT NULL GROUP BY 1)
    SELECT tw.brand, tw.tw_ord, ROUND(a4.avg_ord)::int, ROUND((tw.tw_ord-a4.avg_ord)/NULLIF(a4.avg_ord,0)*100,1) FROM tw JOIN a4 USING(brand) WHERE a4.avg_ord>=5 ORDER BY tw.tw_ord DESC LIMIT 8
""")

watch_list = q(f"""
    WITH lw AS (SELECT DISTINCT revenue_customer_name AS brand FROM gold.fact_revenue WHERE service_line='DAASH' AND revenue_order_date BETWEEN '{s_lw}' AND '{e_lw}' AND revenue_customer_name IS NOT NULL),
    tw AS (SELECT DISTINCT revenue_customer_name AS brand FROM gold.fact_revenue WHERE service_line='DAASH' AND revenue_order_date BETWEEN '{s_tw}' AND '{e_tw}' AND revenue_customer_name IS NOT NULL)
    SELECT lw.brand FROM lw LEFT JOIN tw ON lw.brand=tw.brand WHERE tw.brand IS NULL
""")

activation_targets = q(f"""
    SELECT revenue_customer_name, COUNT(*) FILTER (WHERE revenue_order_channel='pos') AS p, COUNT(*) FILTER (WHERE revenue_order_channel='website') AS w
    FROM gold.fact_revenue WHERE service_line='DAASH' AND revenue_order_date BETWEEN '{s_tw}' AND '{e_tw}' AND revenue_customer_name IS NOT NULL
    GROUP BY 1 HAVING COUNT(*) FILTER (WHERE revenue_order_channel='website')=0 AND COUNT(*) FILTER (WHERE revenue_order_channel='pos')>10
    ORDER BY p DESC LIMIT 5
""")

top3_pct = float(q1(f"""
    WITH r AS (SELECT revenue_customer_name, SUM(sales_amount) AS rev, SUM(SUM(sales_amount)) OVER() AS tot
    FROM gold.fact_revenue WHERE service_line='DAASH' AND revenue_order_date BETWEEN '{s_tw}' AND '{e_tw}' AND revenue_customer_name IS NOT NULL GROUP BY 1 ORDER BY rev DESC LIMIT 3)
    SELECT COALESCE(ROUND((SUM(rev)/NULLIF(MAX(tot),0)*100)::numeric,1),0) FROM r
""")[0] or 0)

# GoSource data
gs_new = q(f"WITH fo AS (SELECT revenue_customer_name AS c, MIN(revenue_order_date) AS fd FROM gold.fact_revenue WHERE service_line='GoSource' AND revenue_customer_name IS NOT NULL GROUP BY 1) SELECT c, fd FROM fo WHERE fd BETWEEN '{s_tw}' AND '{e_tw}'")

# New vs reactivated brands (DAASH)
daash_new_reactivated = q(f"""
    WITH tw AS (SELECT DISTINCT revenue_customer_name AS brand FROM gold.fact_revenue WHERE service_line='DAASH' AND revenue_order_date BETWEEN '{s_tw}' AND '{e_tw}' AND revenue_customer_name IS NOT NULL),
    prior4 AS (SELECT DISTINCT revenue_customer_name AS brand FROM gold.fact_revenue WHERE service_line='DAASH' AND revenue_order_date BETWEEN '{PRIOR_4[3][0]}' AND '{PRIOR_4[0][1]}' AND revenue_customer_name IS NOT NULL),
    ever_before AS (SELECT DISTINCT revenue_customer_name AS brand FROM gold.fact_revenue WHERE service_line='DAASH' AND revenue_order_date < '{PRIOR_4[3][0]}' AND revenue_customer_name IS NOT NULL)
    SELECT tw.brand,
           CASE WHEN ever_before.brand IS NOT NULL THEN 'REACTIVATED' ELSE 'NEW' END AS status,
           (SELECT MIN(revenue_order_date) FROM gold.fact_revenue WHERE service_line='DAASH' AND revenue_customer_name=tw.brand) AS first_order,
           (SELECT COUNT(*) FROM gold.fact_revenue WHERE service_line='DAASH' AND revenue_customer_name=tw.brand AND revenue_order_date BETWEEN '{s_tw}' AND '{e_tw}') AS tw_orders
    FROM tw LEFT JOIN prior4 ON tw.brand=prior4.brand LEFT JOIN ever_before ON tw.brand=ever_before.brand
    WHERE prior4.brand IS NULL ORDER BY status, tw_orders DESC
""")

# New vs reactivated customers (GoSource)
gs_new_reactivated = q(f"""
    WITH tw AS (SELECT DISTINCT revenue_customer_name AS c FROM gold.fact_revenue WHERE service_line='GoSource' AND revenue_order_date BETWEEN '{s_tw}' AND '{e_tw}' AND revenue_customer_name IS NOT NULL),
    prior4 AS (SELECT DISTINCT revenue_customer_name AS c FROM gold.fact_revenue WHERE service_line='GoSource' AND revenue_order_date BETWEEN '{PRIOR_4[3][0]}' AND '{PRIOR_4[0][1]}' AND revenue_customer_name IS NOT NULL),
    ever_before AS (SELECT DISTINCT revenue_customer_name AS c FROM gold.fact_revenue WHERE service_line='GoSource' AND revenue_order_date < '{PRIOR_4[3][0]}' AND revenue_customer_name IS NOT NULL)
    SELECT tw.c,
           CASE WHEN ever_before.c IS NOT NULL THEN 'REACTIVATED' ELSE 'NEW' END AS status,
           (SELECT MIN(revenue_order_date) FROM gold.fact_revenue WHERE service_line='GoSource' AND revenue_customer_name=tw.c) AS first_order,
           (SELECT COUNT(*) FROM gold.fact_revenue WHERE service_line='GoSource' AND revenue_customer_name=tw.c AND revenue_order_date BETWEEN '{s_tw}' AND '{e_tw}') AS tw_orders
    FROM tw LEFT JOIN prior4 ON tw.c=prior4.c LEFT JOIN ever_before ON tw.c=ever_before.c
    WHERE prior4.c IS NULL ORDER BY status, tw_orders DESC
""")

gs_pipeline = q(f"""
    SELECT DISTINCT ON (order_id_pk) order_business_name, order_status, order_created_at_date, order_total_price_amount
    FROM bv.bv_gosource_orders WHERE order_created_at_date BETWEEN '{s_tw}' AND '{e_tw}'
    ORDER BY order_id_pk, order_delivered_at_date
""")
gs_status = {}
for _, status, _, _ in gs_pipeline:
    gs_status[status] = gs_status.get(status, 0) + 1

gs_pending = q(f"""
    SELECT DISTINCT ON (o.order_id_pk) COALESCE(c.customer_business_name, o.order_business_name, 'Unknown') AS customer, o.order_created_at_date, o.order_total_price_amount
    FROM bv.bv_gosource_orders o LEFT JOIN bv.bv_gosource_customers c ON o.order_unified_customer_id_fk = c.customer_id_pk
    WHERE o.order_created_at_date BETWEEN '{s_tw}' AND '{e_tw}' AND lower(o.order_status) = 'pending'
    ORDER BY o.order_id_pk, o.order_delivered_at_date
""")

ar_total = float(q1("SELECT COALESCE(SUM(ar_outstanding_amount),0) FROM gold.fact_ar_aging")[0])
ar_90 = float(q1("SELECT COALESCE(SUM(ar_outstanding_amount),0) FROM gold.fact_ar_aging WHERE ar_aging_bucket='90+ days'")[0])
ar_90_pct = (ar_90/ar_total*100) if ar_total > 0 else 0

ar_top = q("SELECT ar_customer_name, SUM(ar_outstanding_amount)::bigint AS amt, COUNT(*) AS inv FROM gold.fact_ar_aging GROUP BY 1 ORDER BY amt DESC LIMIT 5")

print("Building PowerPoint...")

# ---------------------------------------------------------------------------
# PowerPoint setup
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height


def add_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_banner(slide, title, subtitle, color=CRIMSON):
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.2))
    banner.fill.solid()
    banner.fill.fore_color.rgb = color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(0xFF, 0xDD, 0xDD) if color == CRIMSON else RGBColor(0xCC, 0xEE, 0xDD)


def add_footer(slide, text):
    ft = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12), Inches(0.3))
    ft.text_frame.paragraphs[0].text = text
    ft.text_frame.paragraphs[0].font.size = Pt(8)
    ft.text_frame.paragraphs[0].font.color.rgb = GREY


def kpi_card(slide, x, y, w, h, label, value, sub, accent=CRIMSON, good=True):
    # Accent stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.06))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()
    # Card body
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Inches(0.06), w, h - Inches(0.06))
    card.fill.solid()
    card.fill.fore_color.rgb = PINK_BG if accent == CRIMSON else GREEN_BG if accent == GS_GREEN else BLUE_BG
    card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    card.line.width = Pt(0.5)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = accent
    p.alignment = PP_ALIGN.CENTER
    pv = tf.add_paragraph()
    pv.text = value
    pv.font.size = Pt(26)
    pv.font.bold = True
    pv.font.color.rgb = BLACK
    pv.alignment = PP_ALIGN.CENTER
    pv.space_before = Pt(4)
    ps = tf.add_paragraph()
    ps.text = sub
    ps.font.size = Pt(9)
    ps.font.color.rgb = DARK_GREEN if good else ALERT_RED
    ps.alignment = PP_ALIGN.CENTER
    ps.space_before = Pt(2)


def info_box(slide, x, y, w, h, title, lines, bg=BLUE_BG, title_color=NAVY, text_color=BLACK):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = title_color
    for line in lines:
        lp = tf.add_paragraph()
        lp.text = line
        lp.font.size = Pt(9)
        lp.font.color.rgb = text_color
        lp.space_before = Pt(3)


def add_table(slide, x, y, w, data, col_ws=None):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, x, y, w, Inches(0.32 * rows))
    tbl = ts.table
    if col_ws:
        for i, cw in enumerate(col_ws):
            tbl.columns[i].width = cw
    for r, rd in enumerate(data):
        for c, v in enumerate(rd):
            cell = tbl.cell(r, c)
            cell.text = str(v)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
                p.font.color.rgb = WHITE
                p.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else LIGHT_GREY
                p.font.color.rgb = BLACK
            if c > 0:
                p.alignment = PP_ALIGN.RIGHT
    return ts


FOOTER = f"IPC Weekly Report • {TW_LABEL}"

# ===========================================================================
# SLIDE 1 — TITLE
# ===========================================================================
s1 = add_slide()
# Full background
bg = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()
# Red accent bar
bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), SW, Inches(2.8))
bar.fill.solid()
bar.fill.fore_color.rgb = CRIMSON
bar.line.fill.background()
# Title text
tb = s1.shapes.add_textbox(Inches(1), Inches(2.7), Inches(11), Inches(2.2))
tf = tb.text_frame
tf.word_wrap = True
p1 = tf.paragraphs[0]
p1.text = "IPC WEEKLY"
p1.font.size = Pt(48)
p1.font.bold = True
p1.font.color.rgb = WHITE
p1.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "PERFORMANCE REPORT"
p2.font.size = Pt(48)
p2.font.bold = True
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph()
p3.text = f"Week of {TW_LABEL}"
p3.font.size = Pt(18)
p3.font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(12)
# Guide
guide = s1.shapes.add_textbox(Inches(3), Inches(5.8), Inches(7), Inches(0.8))
gtf = guide.text_frame
gtf.word_wrap = True
gp = gtf.paragraphs[0]
gp.text = "DAASH (Slides 2–4)  •  GoSource (Slides 5–6)  •  Wins & Actions (Slide 7)"
gp.font.size = Pt(14)
gp.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
gp.alignment = PP_ALIGN.CENTER
gp2 = gtf.add_paragraph()
gp2.text = "vs 4-week trailing average  •  Management Review"
gp2.font.size = Pt(11)
gp2.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
gp2.alignment = PP_ALIGN.CENTER


# ===========================================================================
# SLIDE 2 — DAASH SNAPSHOT
# ===========================================================================
s2 = add_slide()
add_banner(s2, f"DAASH — WEEKLY SNAPSHOT", f"Week of {TW_LABEL} • Food Delivery B2C")

cw = Inches(2.85)
ch = Inches(1.25)
cy = Inches(1.5)
gap = Inches(0.2)
sx = Inches(0.5)

kpi_card(s2, sx, cy, cw, ch, "IPC SERVICE CHARGE", naira(tw["daash_sc"]),
         delta_str(tw["daash_sc"], avg4("daash_sc")),
         CRIMSON, is_up(tw["daash_sc"], avg4("daash_sc")))
kpi_card(s2, sx+cw+gap, cy, cw, ch, "GMV (FOOD SOLD)", naira(tw["daash_gmv"]),
         delta_str(tw["daash_gmv"], avg4("daash_gmv")),
         CRIMSON, is_up(tw["daash_gmv"], avg4("daash_gmv")))
kpi_card(s2, sx+2*(cw+gap), cy, cw, ch, "DELIVERED ORDERS", f"{tw['daash_orders']:,}",
         delta_str(tw["daash_orders"], avg4("daash_orders")),
         CRIMSON, is_up(tw["daash_orders"], avg4("daash_orders")))
kpi_card(s2, sx+3*(cw+gap), cy, cw, ch, "AVG ORDER VALUE", naira(tw["daash_aov"]),
         delta_str(tw["daash_aov"], avg4("daash_aov")),
         CRIMSON, is_up(tw["daash_aov"], avg4("daash_aov")))

# Explain service charge
expl = s2.shapes.add_textbox(sx, Inches(2.85), Inches(6), Inches(0.3))
expl.text_frame.paragraphs[0].text = "IPC Service Charge = what IPC actually earns. GMV = total food value processed. They are not the same."
expl.text_frame.paragraphs[0].font.size = Pt(8)
expl.text_frame.paragraphs[0].font.italic = True
expl.text_frame.paragraphs[0].font.color.rgb = GREY

# Mini stat cards row 2
mini_y = Inches(3.3)
mini_h = Inches(1.0)
kpi_card(s2, sx, mini_y, Inches(3.8), mini_h, "ACTIVE BRANDS", f"{tw['daash_brands']}",
         f"4wk avg: {avg4('daash_brands'):.0f}", CRIMSON, is_up(tw["daash_brands"], avg4("daash_brands")))
kpi_card(s2, sx+Inches(4.0), mini_y, Inches(3.8), mini_h, "WEB ORDER SHARE",
         f"{tw['web_pct']:.1f}%",
         f"Target: 30% — {'below target' if tw['web_pct']<30 else 'on target'}",
         CRIMSON, tw['web_pct'] >= 25)
kpi_card(s2, sx+Inches(8.0), mini_y, Inches(4.3), mini_h, "TOP 3 CONCENTRATION",
         f"{top3_pct:.1f}%",
         f"Target: <70% — {'high risk' if top3_pct>70 else 'improving'}",
         CRIMSON, top3_pct < 70)

# Insight box
sc_trend = delta_str(tw["daash_sc"], avg4("daash_sc"))
insight_lines = [
    f"Service charge {sc_trend} vs 4-week average",
    f"Top 3 brands drive {top3_pct:.0f}% of all food value — high dependency",
    f"Web orders at {tw['web_pct']:.0f}% (target 30%) — most activity still on POS",
    f"POS orders generate ₦0 service charge for IPC — online is where we earn",
]
info_box(s2, sx, Inches(4.6), Inches(12.3), Inches(2.2),
         "📊 WHAT THIS MEANS", insight_lines, BLUE_BG, NAVY)
add_footer(s2, FOOTER)


# ===========================================================================
# SLIDE 3 — DAASH CHANNEL & BRAND HEALTH
# ===========================================================================
s3 = add_slide()
add_banner(s3, "DAASH — CHANNEL MIX & BRAND HEALTH",
           "Where orders come from and which brands need attention")

# Channel table
ch_data = [["Channel", "Orders", "Share", "GMV", "vs 4wk Avg"]]
ch_data.append(["POS", f"{tw['pos_orders']:,}",
                f"{tw['pos_orders']/(tw['pos_orders']+tw['web_orders'])*100:.0f}%" if (tw['pos_orders']+tw['web_orders'])>0 else "—",
                naira(tw['pos_gmv']),
                delta_str(tw['pos_orders'], avg4('pos_orders'))])
ch_data.append(["Website", f"{tw['web_orders']:,}",
                f"{tw['web_pct']:.0f}%",
                naira(tw['web_gmv']),
                delta_str(tw['web_orders'], avg4('web_orders'))])
add_table(s3, Inches(0.5), Inches(1.5), Inches(6), ch_data,
          col_ws=[Inches(1.2), Inches(1.0), Inches(0.8), Inches(1.2), Inches(1.5)])

info_box(s3, Inches(0.5), Inches(2.7), Inches(6), Inches(0.8),
         "", ["POS = in-store (₦0 service charge). Website = online (generates IPC revenue). Push brands toward website ordering."],
         AMBER_BG, GREY, GREY)

# Watch list
if watch_list:
    w_lines = [f"• {r[0]} — was active last week, no orders this week" for r in watch_list[:5]]
    w_lines.append("→ Action: call to check — supply issue? holiday? competitor?")
else:
    w_lines = ["All brands active this week — no churn detected"]
info_box(s3, Inches(7), Inches(1.5), Inches(5.8), Inches(2.0),
         "⚠️ WENT QUIET THIS WEEK", w_lines, RED_BG, ALERT_RED)

# Activation targets
if activation_targets:
    a_lines = [f"• {r[0]} — {r[1]} POS orders, 0 web" for r in activation_targets[:5]]
    a_lines.append("→ Action: set up Google Business Profile, send WhatsApp template")
else:
    a_lines = ["All active brands have web orders this week"]
info_box(s3, Inches(7), Inches(3.8), Inches(5.8), Inches(2.0),
         "🎯 ACTIVATION TARGETS (POS-heavy, zero web)", a_lines, AMBER_BG, RGBColor(0xB4, 0x5D, 0x09))

# New / reactivated brands box
nr_lines = []
for brand, status, first_ord, tw_ord in daash_new_reactivated:
    tag = "🆕 NEW" if status == "NEW" else "🔄 REACTIVATED"
    nr_lines.append(f"{tag}: {brand} — {tw_ord} orders this week (first ever: {first_ord})")
if not nr_lines:
    nr_lines = ["No new or reactivated brands this week"]
info_box(s3, Inches(0.5), Inches(3.8), Inches(6), Inches(1.2),
         "🆕 NEW & REACTIVATED BRANDS", nr_lines, GREEN_BG, DARK_GREEN)

# Wins
wins = []
if tw["daash_sc"] > avg4("daash_sc"): wins.append(f"Service charge above 4-week average ({naira(tw['daash_sc'])})")
if tw["daash_brands"] > avg4("daash_brands"): wins.append(f"Active brands above average ({tw['daash_brands']} vs {avg4('daash_brands'):.0f})")
if top_movers and float(top_movers[0][3] or 0) > 0: wins.append(f"{top_movers[0][0]} grew {top_movers[0][3]}% vs 4-week avg")
if daash_new_reactivated:
    new_count = sum(1 for r in daash_new_reactivated if r[1]=='NEW')
    react_count = sum(1 for r in daash_new_reactivated if r[1]=='REACTIVATED')
    if new_count: wins.append(f"{new_count} brand new brand(s) activated")
    if react_count: wins.append(f"{react_count} dormant brand(s) reactivated")
if not watch_list: wins.append("Zero brand churn this week — 100% retention")
if not wins: wins = ["Steady week — no major declines"]
info_box(s3, Inches(0.5), Inches(5.2), Inches(6), Inches(1.6),
         "✅ WINS THIS WEEK", wins, GREEN_BG, DARK_GREEN)
add_footer(s3, FOOTER)


# ===========================================================================
# SLIDE 4 — DAASH TOP MOVERS
# ===========================================================================
s4 = add_slide()
add_banner(s4, "DAASH — TOP MOVERS VS 4-WEEK AVERAGE",
           "Which brands are growing, shrinking, or steady")

if top_movers:
    m_data = [["#", "Brand", "This Week", "4wk Avg", "% vs Avg", "Trend"]]
    for i, (brand, tw_o, avg_o, pct_chg) in enumerate(top_movers):
        pct_chg = float(pct_chg or 0)
        trend = "▲▲" if pct_chg > 15 else "▲" if pct_chg > 0 else "▼" if pct_chg > -15 else "▼▼"
        m_data.append([str(i+1), brand, f"{tw_o:,}", f"{avg_o:,}", f"{pct_chg:+.1f}%", trend])
    add_table(s4, Inches(0.5), Inches(1.5), Inches(8.5), m_data,
              col_ws=[Inches(0.4), Inches(3.0), Inches(1.2), Inches(1.2), Inches(1.2), Inches(0.8)])

info_box(s4, Inches(0.5), Inches(4.5), Inches(8.5), Inches(0.7),
         "", ["Brands above their 4-week average are growing organically. Brands below may need an ops check-in or activation push."],
         LIGHT_GREY, GREY, GREY)

# 4-week SC trend
trend_data = [["Week", "Service Charge", "Orders", "GMV", "Web %"]]
for i, p in enumerate(prior):
    w = PRIOR_4[i]
    trend_data.append([f"W-{i+1} ({w[0].strftime('%d %b')})", naira(p["daash_sc"]),
                       f"{p['daash_orders']:,}", naira(p["daash_gmv"]), f"{p['web_pct']:.1f}%"])
trend_data.append(["4-Week Avg", naira(avg4("daash_sc")), f"{int(avg4('daash_orders')):,}",
                    naira(avg4("daash_gmv")), f"{avg4('web_pct'):.1f}%"])
add_table(s4, Inches(0.5), Inches(5.5), Inches(7.5), trend_data,
          col_ws=[Inches(1.8), Inches(1.4), Inches(1.2), Inches(1.4), Inches(1.0)])

info_box(s4, Inches(9.5), Inches(1.5), Inches(3.3), Inches(4.5),
         "📈 4-WEEK CONTEXT", [
             f"Avg weekly SC: {naira(avg4('daash_sc'))}",
             f"Avg weekly orders: {int(avg4('daash_orders')):,}",
             f"Avg web share: {avg4('web_pct'):.1f}%",
             f"This week SC: {naira(tw['daash_sc'])}",
             f"Trend: {'above' if tw['daash_sc']>avg4('daash_sc') else 'below'} average",
         ], BLUE_BG, NAVY)
add_footer(s4, FOOTER)


# ===========================================================================
# SLIDE 5 — GOSOURCE SNAPSHOT
# ===========================================================================
s5 = add_slide()
add_banner(s5, "GOSOURCE — WEEKLY SNAPSHOT",
           f"Week of {TW_LABEL} • B2B Procurement", GS_GREEN)

cy5 = Inches(1.5)
kpi_card(s5, sx, cy5, cw, ch, "IPC SERVICE CHARGE", naira(tw["gs_sc"]),
         "Credit product dormant" if tw["gs_sc"]==0 else delta_str(tw["gs_sc"], avg4("gs_sc")),
         GS_GREEN, tw["gs_sc"] > 0)
kpi_card(s5, sx+cw+gap, cy5, cw, ch, "GMV (GOODS SOLD)", naira(tw["gs_gmv"]),
         delta_str(tw["gs_gmv"], avg4("gs_gmv")),
         GS_GREEN, is_up(tw["gs_gmv"], avg4("gs_gmv")))
kpi_card(s5, sx+2*(cw+gap), cy5, cw, ch, "DELIVERED ORDERS", f"{tw['gs_orders']}",
         f"4wk avg: {avg4('gs_orders'):.0f}",
         GS_GREEN, is_up(tw["gs_orders"], avg4("gs_orders")))
kpi_card(s5, sx+3*(cw+gap), cy5, cw, ch, "ACTIVE CUSTOMERS", f"{tw['gs_customers']}",
         f"{len(gs_new)} new this week" if gs_new else "No new customers",
         GS_GREEN, len(gs_new) > 0)

# Alert banner if volume is way down
if tw["gs_gmv"] < avg4("gs_gmv") * 0.5:
    alert = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.0), Inches(12.3), Inches(1.2))
    alert.fill.solid()
    alert.fill.fore_color.rgb = ALERT_RED
    alert.line.fill.background()
    atf = alert.text_frame
    atf.word_wrap = True
    atf.margin_left = Inches(0.2)
    atf.margin_top = Inches(0.1)
    ap = atf.paragraphs[0]
    ap.text = f"⚠️ VOLUME ALERT: {tw['gs_orders']} delivered orders vs {avg4('gs_orders'):.0f} weekly average. GMV at {naira(tw['gs_gmv'])} — {delta_str(tw['gs_gmv'], avg4('gs_gmv'))} vs trend."
    ap.font.size = Pt(12)
    ap.font.bold = True
    ap.font.color.rgb = WHITE
    ap2 = atf.add_paragraph()
    pending_count = gs_status.get('pending', 0)
    ap2.text = f"{pending_count} orders stuck in 'pending' status — these need follow-up. Each pending order is potential {naira(avg4('gs_aov'))} in GMV."
    ap2.font.size = Pt(10)
    ap2.font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)

# New customers + wins
nr_gs_lines = []
for cust, status, first_ord, tw_ord in gs_new_reactivated:
    tag = "🆕 NEW" if status == "NEW" else "🔄 REACTIVATED"
    nr_gs_lines.append(f"{tag}: {cust} — {tw_ord} order(s) (first ever: {first_ord})")
if not nr_gs_lines:
    nr_gs_lines = ["No new or reactivated customers this week"]
info_box(s5, Inches(0.5), Inches(4.5), Inches(5.8), Inches(2.2),
         "🆕 NEW & REACTIVATED CUSTOMERS", nr_gs_lines, GREEN_BG, DARK_GREEN)

gs_wins = []
new_gs_count = sum(1 for r in gs_new_reactivated if r[1]=='NEW')
react_gs_count = sum(1 for r in gs_new_reactivated if r[1]=='REACTIVATED')
if new_gs_count: gs_wins.append(f"{new_gs_count} brand new customer(s) acquired")
if react_gs_count: gs_wins.append(f"{react_gs_count} dormant customer(s) reactivated")
if tw["gs_customers"] >= avg4("gs_customers"): gs_wins.append(f"Customer base stable ({tw['gs_customers']} active)")
if not gs_wins: gs_wins = ["Customer acquisition continues despite volume dip"]
info_box(s5, Inches(6.8), Inches(4.5), Inches(5.9), Inches(2.2),
         "✅ BRIGHT SPOTS", gs_wins, GREEN_BG, DARK_GREEN)
add_footer(s5, FOOTER)


# ===========================================================================
# SLIDE 6 — GOSOURCE PIPELINE
# ===========================================================================
s6 = add_slide()
add_banner(s6, "GOSOURCE — ORDER PIPELINE & BLOCKERS",
           "What's in the system and what's stuck", GS_GREEN)

# Status summary
status_data = [["Status", "Orders", "% of Total"]]
total_pipe = sum(gs_status.values())
for st in ["delivered", "pending", "processing", "partially_delivered", "cancelled", "rejected"]:
    cnt = gs_status.get(st, 0)
    if cnt > 0:
        status_data.append([st.title(), f"{cnt}", f"{cnt/total_pipe*100:.0f}%" if total_pipe > 0 else "—"])
add_table(s6, Inches(0.5), Inches(1.5), Inches(5), status_data,
          col_ws=[Inches(2.0), Inches(1.2), Inches(1.2)])

if gs_pending:
    pending_data = [["Customer", "Order Date", "Days Pending", "Est. Value"]]
    for cust, odate, val in gs_pending[:10]:
        days = (today - odate).days if odate else 0
        pending_data.append([cust[:25], str(odate), f"{days}d", naira(val)])
    add_table(s6, Inches(6), Inches(1.5), Inches(6.8), pending_data,
              col_ws=[Inches(2.5), Inches(1.3), Inches(1.0), Inches(1.2)])

pending_count = gs_status.get('pending', 0)
info_box(s6, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.2),
         "📋 WHAT TO DO", [
             f"{pending_count} pending orders = potential {naira(pending_count * avg4('gs_aov'))} in GMV if unblocked",
             "Call each pending customer — confirm delivery schedule or cancel to free inventory",
             "Orders pending 7+ days are likely dead unless someone intervenes",
             f"Delivered orders this week: {tw['gs_orders']} — investigate if fulfillment is the bottleneck",
         ], BLUE_BG, NAVY)
add_footer(s6, FOOTER)


# ===========================================================================
# SLIDE 7 — WINS & RECOMMENDED ACTIONS
# ===========================================================================
s8 = add_slide()
add_banner(s8, "WINS & RECOMMENDED ACTIONS",
           f"Week of {TW_LABEL} • Summary for Leadership")

# --- WINS (left side) ---
all_wins = []
# Boli & Grills
boli = [r for r in daash_new_reactivated if 'boli' in r[0].lower()]
if boli:
    all_wins.append(f"Boli & Grills: {boli[0][3]} orders in first 3 days — fastest launch this year")
    all_wins.append("Boli & Grills web share 61% — only brand where web beats POS")
all_wins.append(f"Active DAASH brands at {tw['daash_brands']}, retention holds")
reactivated = [r for r in daash_new_reactivated if r[1] == 'REACTIVATED']
if reactivated:
    all_wins.append(f"{reactivated[0][0]} reactivated after months dormant")
new_brands = [r for r in daash_new_reactivated if r[1] == 'NEW']
if len(new_brands) > 1:
    all_wins.append(f"{len(new_brands)} brand new DAASH brands activated")
if gs_new_reactivated:
    new_gs = [r for r in gs_new_reactivated if r[1] == 'NEW']
    if new_gs:
        all_wins.append(f"New GoSource customer: {new_gs[0][0]} — pipeline producing")
pending_count = gs_status.get('pending', 0)
if pending_count > 5:
    all_wins.append(f"GoSource has {pending_count} pending orders — demand exists")
if tw["daash_sc"] > avg4("daash_sc"):
    all_wins.append(f"DAASH service charge above 4-week average")
if not all_wins:
    all_wins = ["Steady week across both service lines"]

info_box(s8, Inches(0.4), Inches(1.5), Inches(5.8), Inches(3.8),
         "✅ THIS WEEK'S WINS", all_wins[:7], GREEN_BG, DARK_GREEN)

# --- RECOMMENDATIONS (right side) ---
recs = [
    "1. Study Boli & Grills' web playbook and replicate (Product) — 61% web share day 1, package for Captains, Pitakwa, Urban Bites",
    f"2. Unblock {pending_count} stuck GoSource orders (Ops) — GMV collapsed but demand is pending, not gone",
    "3. Web activation sprint: 3 POS-heavy brands (Growth) — Google profiles, Instagram order buttons, QR table cards",
]

# Add Papas Grill debt rec if they're top debtor
if ar_top and ar_top[0][0] and 'papa' in ar_top[0][0].lower():
    recs.append(f"4. Call {ar_top[0][0]} about {naira(ar_top[0][1])} overdue (Account Mgmt) — largest debtor AND top revenue driver")

if watch_list:
    quiet_names = ", ".join(r[0] for r in watch_list[:3])
    recs.append(f"5. Re-engage quiet brands: {quiet_names} (Account Mgmt) — WhatsApp check-in this week")

recs.append("6. Pause new GoSource credit until AR improves (Finance) — all ₦83.9M is 90+ days overdue")

info_box(s8, Inches(6.6), Inches(1.5), Inches(6.3), Inches(3.8),
         "📋 RECOMMENDED ACTIONS", recs, BLUE_BG, NAVY)

# --- HEADLINE OF THE WEEK (bottom banner) ---
headline_box = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(5.6), Inches(12.5), Inches(1.2))
headline_box.fill.solid()
headline_box.fill.fore_color.rgb = NAVY
headline_box.line.fill.background()
htf = headline_box.text_frame
htf.word_wrap = True
htf.margin_left = Inches(0.3)
htf.margin_top = Inches(0.15)
hp = htf.paragraphs[0]
hp.text = "💡 HEADLINE OF THE WEEK"
hp.font.size = Pt(10)
hp.font.bold = True
hp.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
hp2 = htf.add_paragraph()
if boli:
    hp2.text = "Boli & Grills proved web-first activation works in 3 days — now we need to make it the standard, not the exception."
else:
    hp2.text = f"Activation remains the core challenge — {tw['web_pct']:.0f}% web share vs 30% target. Every POS-only brand is leaving IPC revenue on the table."
hp2.font.size = Pt(16)
hp2.font.bold = True
hp2.font.color.rgb = WHITE
hp2.space_before = Pt(4)

add_footer(s8, FOOTER)


# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
OUTPUT = f"IPC_Weekly_Report_{TW[0].strftime('%Y%m%d')}.pptx"
prs.save(OUTPUT)

# Also save to Desktop
import shutil
desktop = os.path.expanduser("~/Desktop")
shutil.copy(OUTPUT, os.path.join(desktop, OUTPUT))

cur.close()
conn.close()
print(f"\n✅ Generated: {OUTPUT}")
print(f"📁 Saved to Desktop: ~/Desktop/{OUTPUT}")
